#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generator for the real-document fixtures used by the file-preview E2E tests.

WHAT THIS IS
------------
The files in ``e2e/fixtures/files/`` are **committed to the repository**. CI never
runs this script, and no test imports it. It exists for two reasons only:

1. *Provenance* — so a reviewer can see exactly how every binary fixture was
   produced and that it contains no real customer data.
2. *Regeneration* — so a fixture can be rebuilt or adjusted deliberately
   instead of being a mystery blob someone once attached to a PR.

Run it manually when a fixture must change::

    python3 e2e/fixtures/generate_fixtures.py            # generate + verify
    python3 e2e/fixtures/generate_fixtures.py --verify   # verify only, no writes

The verification pass re-opens every file with a real parser and asserts the
properties the specs depend on (page counts, row counts, slide aspect, the
350 KB renderer cutoff, "this one must fail to open"). Run it after any edit.

CONTENT POLICY
--------------
This repository is public. Every company name, person, address, bank detail,
sum, e-mail and date below is **invented**: ООО «Ромашка», ИП Иванов И. И.,
«ул. Придуманная», УНП 100000001, ``@example.com`` (an RFC 2606 reserved
domain). Nothing here is derived from a real contract, a real counterparty or
any production system. Keep it that way when editing.

The prose is deliberately Russian business language (the product UI is
Russian), with a little English where it occurs naturally, because Cyrillic
encoding and font handling are part of what the preview tests exercise.

DEPENDENCIES
------------
Standard library plus: ``python-docx``, ``openpyxl``, ``python-pptx``,
``reportlab``, ``pikepdf``, ``Pillow``.

``msoffcrypto-tool`` is **optional** and needed only for ``locked.docx`` (a
genuinely encrypted Office file). Without it that single fixture is skipped
with a loud warning — the script never fabricates a fake "encrypted" file.

DETERMINISM
-----------
Two runs on the same machine produce byte-identical files, so a regeneration
that changes nothing shows up as an empty diff:

* all document metadata dates are pinned to ``FIXED_DT`` (2024-01-15 09:00);
* every OOXML zip is rewritten with fixed entry timestamps, fixed permissions
  and a fixed compression level (``normalize_zip``);
* reportlab runs with ``invariant=1`` (no wall-clock CreationDate, stable
  document ID);
* all "random-looking" bytes come from a SHA-256 keystream with a fixed seed,
  never from ``random`` or ``os.urandom``.

Two documented exceptions, both unavoidable: ``locked.pdf`` and ``locked.docx``
embed cryptographic salts, so their bytes differ on every run even though their
content does not. Do not re-commit them unless their content actually changed.

The embedded raster images are JPEG-encoded by Pillow, so byte-identical output
also assumes the same Pillow/libjpeg build; the fixtures stay valid either way.
"""

from __future__ import annotations

import argparse
import csv
import hashlib
import io
import os
import re
import sys
import zipfile
from datetime import datetime
from pathlib import Path

# --------------------------------------------------------------------------
# Constants
# --------------------------------------------------------------------------

HERE = Path(__file__).resolve().parent
OUT_DIR = HERE / "files"

#: Every metadata date in every generated document. Never "now".
FIXED_DT = datetime(2024, 1, 15, 9, 0, 0)
#: Zip entry timestamp (zip stores local time as a 6-tuple).
ZIP_DT = (2024, 1, 15, 9, 0, 0)
#: Fixed mode bits for zip entries, so umask cannot leak into the output.
ZIP_ATTR = 0o644 << 16

DOCX_PASSWORD = "secret123"  # documented in files/README.md, fixtures only
PDF_PASSWORD = "secret123"

#: Backend constants this fixture set is built against
#: (packages/api/src/files/documents/html.ts).
MAX_DOCX_CDN_BINARY_BYTES = 350 * 1024
SPREADSHEET_MAX_ROWS_PER_SHEET = 5_000

#: contract-heavy.docx must clear the CDN cutoff with real margin.
HEAVY_TARGET_BYTES = 420 * 1024

#: contract-long.docx volume knobs. Text compresses hard inside a docx, so
#: reaching the 60-200 KB band takes a genuinely long agreement: each clause
#: pool is cycled LONG_CLAUSE_REPEATS times and LONG_APPENDIX_COUNT appendices
#: follow, every one of them starting on a fresh page.
LONG_CLAUSE_REPEATS = 6
LONG_APPENDIX_COUNT = 12
#: Rows in the «Смета» appendix of contract-long.docx.
LONG_ESTIMATE_ROWS = 220

# Fonts that contain Cyrillic, in preference order: open fonts first, then the
# macOS system fonts (fsType=8, "editable embedding" — subsetting into a PDF is
# permitted by the vendor's own embedding bits).
FONT_CANDIDATES = (
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/TTF/DejaVuSans.ttf",
    "/usr/share/fonts/dejavu/DejaVuSans.ttf",
    "/Library/Fonts/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSerif-Regular.ttf",
    "/System/Library/Fonts/Supplemental/Times New Roman.ttf",
    "/System/Library/Fonts/Supplemental/Arial.ttf",
    "/Library/Fonts/Arial Unicode.ttf",
)
BOLD_SUFFIX_MAP = {
    "DejaVuSans.ttf": "DejaVuSans-Bold.ttf",
    "LiberationSerif-Regular.ttf": "LiberationSerif-Bold.ttf",
    "Times New Roman.ttf": "Times New Roman Bold.ttf",
    "Arial.ttf": "Arial Bold.ttf",
}


# --------------------------------------------------------------------------
# Deterministic helpers
# --------------------------------------------------------------------------


def keystream(seed: str, length: int) -> bytes:
    """Deterministic pseudo-random bytes (SHA-256 in counter mode).

    Used wherever the fixtures need incompressible or opaque bytes. Unlike
    ``random`` this is stable across Python versions and platforms.
    """
    out = bytearray()
    counter = 0
    base = seed.encode("utf-8")
    while len(out) < length:
        out += hashlib.sha256(base + counter.to_bytes(8, "big")).digest()
        counter += 1
    return bytes(out[:length])


def normalize_zip(path: Path, transform=None) -> None:
    """Rewrite a zip in place with fixed timestamps, modes and compression.

    OOXML writers stamp the current time into every zip entry, which would make
    a regenerated fixture differ from the committed one on every run. This
    rewrites the archive preserving entry order (``[Content_Types].xml`` stays
    first, as OPC readers expect).

    ``transform(names, data) -> (names, data)`` may add, drop or edit parts.
    """
    with zipfile.ZipFile(path) as zin:
        names = zin.namelist()
        data = {name: zin.read(name) for name in names}

    if transform is not None:
        names, data = transform(names, data)

    # openpyxl ignores the modified date we set and stamps the wall clock at
    # save time; pin it here so a regeneration is a no-op diff.
    core = "docProps/core.xml"
    if core in data:
        stamp = FIXED_DT.strftime("%Y-%m-%dT%H:%M:%SZ")
        xml = data[core].decode("utf-8")
        xml = re.sub(
            r"(<dcterms:(?:created|modified)[^>]*>)[^<]*(</dcterms:)",
            lambda m: m.group(1) + stamp + m.group(2),
            xml,
        )
        data[core] = xml.encode("utf-8")

    tmp = path.with_name(path.name + ".tmp")
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED, compresslevel=9) as zout:
        for name in names:
            info = zipfile.ZipInfo(name, date_time=ZIP_DT)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = ZIP_ATTR
            info.create_system = 0  # "MS-DOS", not the host OS
            zout.writestr(info, data[name])
    tmp.replace(path)


def slim_docx(names, data):
    """Drop two optional parts python-docx inherits from its template.

    The stock template weighs ~36 KB before a single character of content, and
    two thirds of that is ballast we can shed without touching anything a
    renderer reads:

    * ``word/stylesWithEffects.xml`` — a Word 2007 back-compat mirror of
      styles.xml (~13 KB compressed). Documents produced by anything other than
      Word normally do not have it at all.
    * ``docProps/thumbnail.jpeg`` — the template's preview image (~1.5 KB),
      showing an unrelated blank page.

    Both are removed together with their relationship entries and their
    ``[Content_Types].xml`` declarations, so the package stays fully valid
    (``--verify`` re-checks that every relationship target and every content
    type override still resolves). This is what keeps contract-short.docx
    under the 30 KB target.
    """
    drop = {"word/stylesWithEffects.xml", "docProps/thumbnail.jpeg"}
    names = [n for n in names if n not in drop]
    data = {k: v for k, v in data.items() if k not in drop}

    def strip_patterns(part: str, patterns) -> None:
        if part not in data:
            return
        xml = data[part].decode("utf-8")
        for pattern in patterns:
            xml = re.sub(pattern, "", xml)
        data[part] = xml.encode("utf-8")

    strip_patterns(
        "[Content_Types].xml",
        [
            r'<Override PartName="/word/stylesWithEffects\.xml"[^>]*/>',
            r'<Override PartName="/docProps/thumbnail\.jpeg"[^>]*/>',
        ],
    )
    strip_patterns(
        "word/_rels/document.xml.rels",
        [r"<Relationship[^>]*stylesWithEffects[^>]*/>"],
    )
    strip_patterns("_rels/.rels", [r"<Relationship[^>]*thumbnail[^>]*/>"])
    return names, data


def find_font() -> tuple[str, str]:
    """Return (regular, bold) paths of a TTF that actually has Cyrillic glyphs.

    The check renders «А» and compares it with a private-use codepoint: if the
    bitmaps match, the font is drawing .notdef and has no Cyrillic (Bitstream
    Vera, bundled with reportlab, fails exactly this way).
    """
    from PIL import Image, ImageDraw, ImageFont

    def render(font, char: str) -> bytes:
        """Rasterise one character; identical output means the same glyph."""
        img = Image.new("L", (48, 48), 0)
        ImageDraw.Draw(img).text((4, 4), char, font=font, fill=255)
        return img.tobytes()

    for path in FONT_CANDIDATES:
        if not os.path.exists(path):
            continue
        try:
            font = ImageFont.truetype(path, 32)
            cyrillic = render(font, "\u0410")  # CYRILLIC CAPITAL LETTER A
            notdef = render(font, "\ue000")  # private use: always unmapped
        except Exception:
            continue
        if cyrillic != notdef:
            bold = BOLD_SUFFIX_MAP.get(os.path.basename(path))
            bold_path = os.path.join(os.path.dirname(path), bold) if bold else None
            if not (bold_path and os.path.exists(bold_path)):
                bold_path = path
            return path, bold_path
    raise SystemExit(
        "No Cyrillic-capable TTF found. Install DejaVu Sans or edit FONT_CANDIDATES."
    )


# --------------------------------------------------------------------------
# Invented business content (see CONTENT POLICY in the module docstring)
# --------------------------------------------------------------------------

EXECUTOR = "ООО «Ромашка»"
CUSTOMER = "ИП Иванов Иван Иванович"
THIRD_PARTIES = (
    "ООО «Василёк»",
    "ЗАО «Тестовый Мост»",
    "ООО «ТехноСфера-Плюс»",
    "ОДО «Северный Ветер»",
    "ООО «Гранит-Строй»",
    "ООО «Первая Вымышленная Компания»",
)

REQUISITES = (
    ("Организация", EXECUTOR),
    ("Адрес", "220000, г. Минск, ул. Придуманная, д. 12, оф. 305"),
    ("УНП", "100000001 (вымышленный)"),
    ("Банк", "ЗАО «Банк Примерный»"),
    ("Расчётный счёт", "BY00 TEST 3012 0000 0000 0000 0000"),
    ("Телефон", "+375 (17) 000-00-00"),
    ("E-mail", "contracts@example.com"),
)
CUSTOMER_REQUISITES = {
    "Организация": CUSTOMER,
    "Адрес": "220002, г. Минск, ул. Образцовая, д. 3, кв. 18",
    "УНП": "200000002 (вымышленный)",
    "Банк": "ОАО «Банк Условный»",
    "Расчётный счёт": "BY00 TEST 3012 1111 1111 1111 1111",
    "Телефон": "+375 (29) 000-00-00",
    "E-mail": "ivanov@example.com",
}

#: Clause templates. Placeholders are filled with index-derived values, so a
#: long document repeats boilerplate the way real contracts do, but never
#: renders two identical sentences in a row.
CLAUSES = {
    "Предмет договора": [
        "Исполнитель обязуется по заданию Заказчика оказать услуги по {topic}, "
        "а Заказчик обязуется принять результат оказанных услуг и оплатить его "
        "в порядке и на условиях, предусмотренных настоящим Договором.",
        "Перечень, объём и содержание услуг определяются Техническим заданием "
        "(Приложение № {n1} к настоящему Договору), являющимся его "
        "неотъемлемой частью.",
        "Услуги оказываются по адресу: {addr}, если Сторонами письменно не "
        "согласовано иное место оказания услуг.",
        "Срок оказания услуг — {days} календарных дней с даты подписания "
        "Сторонами настоящего Договора, если иное не установлено календарным "
        "планом работ.",
        "Результатом оказанных услуг признаётся комплект документов и "
        "материалов, передаваемых Заказчику по акту сдачи-приёмки в двух "
        "экземплярах, имеющих равную юридическую силу.",
        "Исполнитель вправе привлекать к оказанию услуг третьих лиц, оставаясь "
        "ответственным перед Заказчиком за их действия как за свои собственные.",
    ],
    "Права и обязанности Сторон": [
        "Исполнитель обязуется оказать услуги надлежащего качества в объёме и "
        "в сроки, предусмотренные настоящим Договором и Приложением № {n1}.",
        "Исполнитель обязуется письменно уведомить Заказчика о невозможности "
        "оказания услуг в согласованный срок не позднее чем за {days2} рабочих "
        "дней до наступления такого срока с указанием причин.",
        "Заказчик обязуется передать Исполнителю исходные данные и документы, "
        "необходимые для оказания услуг, в течение {days2} рабочих дней с даты "
        "получения соответствующего письменного запроса.",
        "Заказчик обязуется обеспечить представителям Исполнителя доступ в "
        "помещения по адресу {addr} в рабочие дни с 09:00 до 18:00.",
        "Заказчик вправе во всякое время проверять ход и качество оказываемых "
        "услуг, не вмешиваясь в оперативно-хозяйственную деятельность "
        "Исполнителя.",
        "Исполнитель вправе приостановить оказание услуг с письменным "
        "уведомлением Заказчика, если непредоставление исходных данных делает "
        "дальнейшее оказание услуг невозможным.",
        "Стороны назначают уполномоченных представителей, чьи контактные "
        "данные указаны в разделе «Реквизиты Сторон»; замена представителя "
        "производится письменным уведомлением.",
    ],
    "Стоимость услуг и порядок расчётов": [
        "Стоимость услуг по настоящему Договору составляет {sum} белорусских "
        "рублей, включая налог на добавленную стоимость по ставке {vat}%.",
        "Заказчик перечисляет предварительную оплату в размере {pct}% от "
        "стоимости услуг в течение {days2} банковских дней с даты выставления "
        "счёта Исполнителем.",
        "Окончательный расчёт производится в течение {days3} банковских дней "
        "с даты подписания Сторонами акта сдачи-приёмки оказанных услуг.",
        "Обязательство Заказчика по оплате считается исполненным с даты "
        "зачисления денежных средств на расчётный счёт Исполнителя.",
        "Стоимость услуг может быть изменена только по соглашению Сторон, "
        "оформленному дополнительным соглашением в письменной форме.",
        "Все расходы, связанные с банковским переводом денежных средств, несёт "
        "Сторона-плательщик.",
        "Стороны ежеквартально проводят сверку взаиморасчётов; акт сверки "
        "подписывается в течение {days2} рабочих дней с даты его получения.",
    ],
    "Порядок сдачи и приёмки услуг": [
        "По завершении оказания услуг Исполнитель передаёт Заказчику акт "
        "сдачи-приёмки оказанных услуг и отчёт об оказанных услугах.",
        "Заказчик в течение {days2} рабочих дней с даты получения акта "
        "подписывает его либо направляет Исполнителю мотивированный отказ от "
        "приёмки с перечнем выявленных недостатков.",
        "При неполучении Исполнителем подписанного акта либо мотивированного "
        "отказа в указанный срок услуги считаются оказанными надлежащим "
        "образом и принятыми Заказчиком без замечаний.",
        "Выявленные недостатки устраняются Исполнителем за свой счёт в срок, "
        "не превышающий {days} календарных дней с даты получения "
        "мотивированного отказа.",
        "Услуги могут сдаваться поэтапно; приёмка этапа оформляется отдельным "
        "актом и не освобождает Исполнителя от ответственности за результат в "
        "целом.",
    ],
    "Ответственность Сторон": [
        "За нарушение сроков оказания услуг Исполнитель уплачивает Заказчику "
        "пеню в размере {rate}% от стоимости несвоевременно оказанных услуг за "
        "каждый день просрочки, но не более {cap}% от стоимости услуг.",
        "За нарушение сроков оплаты Заказчик уплачивает Исполнителю пеню в "
        "размере {rate}% от суммы задолженности за каждый день просрочки, но "
        "не более {cap}% от суммы задолженности.",
        "Уплата пени не освобождает Стороны от исполнения обязательств в "
        "натуре и от возмещения причинённых убытков.",
        "Сторона, нарушившая обязательство, возмещает другой Стороне реальный "
        "ущерб; упущенная выгода возмещению не подлежит.",
        "Совокупная ответственность Исполнителя по настоящему Договору "
        "ограничивается суммой {sum} белорусских рублей.",
        "Ответственность за достоверность исходных данных, переданных "
        "Исполнителю, несёт Заказчик.",
    ],
    "Конфиденциальность": [
        "Стороны обязуются не разглашать третьим лицам сведения, ставшие им "
        "известными в связи с исполнением настоящего Договора, в течение всего "
        "срока его действия и {years} лет после его прекращения.",
        "Конфиденциальной признаётся любая информация технического, "
        "коммерческого и организационного характера, переданная в письменной "
        "или электронной форме с пометкой «Конфиденциально».",
        "Не является нарушением раскрытие информации по требованию "
        "государственных органов в случаях, установленных законодательством, "
        "при условии письменного уведомления другой Стороны.",
        "Каждая из Сторон обеспечивает режим хранения конфиденциальной "
        "информации не менее строгий, чем режим, применяемый к собственной "
        "информации аналогичного характера.",
        "Передача конфиденциальной информации привлекаемым третьим лицам "
        "допускается только при условии принятия ими обязательств о "
        "неразглашении, не менее строгих, чем предусмотренные настоящим "
        "разделом.",
    ],
    "Обработка персональных данных": [
        "Стороны обрабатывают персональные данные представителей друг друга "
        "исключительно в целях исполнения настоящего Договора.",
        "Объём обрабатываемых персональных данных ограничивается фамилией, "
        "именем, отчеством, должностью и служебными контактными данными "
        "представителей Сторон.",
        "Стороны принимают правовые, организационные и технические меры для "
        "защиты персональных данных от неправомерного доступа, изменения и "
        "распространения.",
        "Срок хранения персональных данных не превышает {years} лет с даты "
        "прекращения настоящего Договора, если иной срок не установлен "
        "законодательством.",
        "Передача персональных данных третьим лицам без письменного согласия "
        "субъекта персональных данных не допускается, за исключением случаев, "
        "прямо предусмотренных законодательством.",
    ],
    "Обстоятельства непреодолимой силы": [
        "Стороны освобождаются от ответственности за неисполнение обязательств, "
        "если оно явилось следствием обстоятельств непреодолимой силы, "
        "возникших после заключения настоящего Договора.",
        "К обстоятельствам непреодолимой силы относятся стихийные бедствия, "
        "эпидемии, а также акты государственных органов, делающие исполнение "
        "обязательств невозможным (force majeure в терминологии Приложения № 2).",
        "Сторона, для которой создалась невозможность исполнения обязательств, "
        "письменно извещает другую Сторону в течение {days2} рабочих дней с "
        "даты наступления таких обстоятельств.",
        "Если обстоятельства непреодолимой силы действуют более {days} "
        "календарных дней подряд, любая из Сторон вправе отказаться от "
        "исполнения настоящего Договора в одностороннем порядке.",
    ],
    "Порядок разрешения споров": [
        "Все споры и разногласия, возникающие из настоящего Договора, Стороны "
        "разрешают путём переговоров.",
        "До обращения в суд Сторона обязана направить претензию; срок ответа "
        "на претензию составляет {days2} рабочих дней с даты её получения.",
        "Претензия направляется по адресу, указанному в разделе «Реквизиты "
        "Сторон», заказным письмом либо вручается уполномоченному "
        "представителю под роспись.",
        "Не урегулированные в претензионном порядке споры передаются на "
        "рассмотрение суда по месту нахождения ответчика.",
        "К отношениям Сторон применяется материальное право по месту "
        "регистрации Исполнителя.",
    ],
    "Срок действия и порядок изменения Договора": [
        "Настоящий Договор вступает в силу с даты его подписания обеими "
        "Сторонами и действует до {end_date}, а в части взаиморасчётов — до "
        "полного исполнения Сторонами принятых обязательств.",
        "Настоящий Договор считается продлённым на каждый последующий "
        "календарный год, если ни одна из Сторон не заявит о его прекращении "
        "не позднее чем за {days} календарных дней до даты окончания срока.",
        "Все изменения и дополнения к настоящему Договору действительны при "
        "условии их совершения в письменной форме и подписания обеими "
        "Сторонами.",
        "Любая из Сторон вправе отказаться от исполнения настоящего Договора, "
        "письменно уведомив об этом другую Сторону не позднее чем за {days} "
        "календарных дней до предполагаемой даты расторжения.",
        "При досрочном расторжении Стороны производят окончательный расчёт за "
        "фактически оказанные услуги в течение {days3} банковских дней.",
    ],
    "Прочие условия": [
        "Настоящий Договор составлен в двух экземплярах, имеющих равную "
        "юридическую силу, по одному экземпляру для каждой из Сторон.",
        "Документы, переданные посредством электронной почты с адресов, "
        "указанных в разделе «Реквизиты Сторон», признаются имеющими "
        "юридическую силу до момента обмена оригиналами.",
        "Стороны обязуются письменно уведомлять друг друга об изменении "
        "реквизитов в течение {days2} рабочих дней с даты такого изменения.",
        "Уступка прав и перевод долга по настоящему Договору допускаются "
        "только с письменного согласия другой Стороны.",
        "Во всём, что не урегулировано настоящим Договором, Стороны "
        "руководствуются действующим законодательством.",
        "Приложения к настоящему Договору являются его неотъемлемой частью и "
        "подписываются уполномоченными представителями Сторон.",
    ],
}

TOPICS = (
    "техническому сопровождению информационной системы Заказчика",
    "методологическому консультированию по вопросам документооборота",
    "разработке регламентов внутреннего контроля",
    "проведению обучения сотрудников Заказчика",
    "аудиту процессов подготовки отчётности",
    "настройке рабочих мест и учётных записей пользователей",
)
ADDRESSES = (
    "220000, г. Минск, ул. Придуманная, д. 12, оф. 305",
    "220001, г. Минск, пр-т Тестовый, д. 45, каб. 2",
    "223000, г. Вымышленск, ул. Образцовая, д. 7",
)
SERVICE_NAMES = (
    "Обследование процессов",
    "Разработка регламента",
    "Настройка рабочих мест",
    "Обучение сотрудников",
    "Сопровождение (1 месяц)",
    "Подготовка отчётности",
    "Консультации по запросу",
    "Аудит учётных записей",
)


def clause_params(i: int) -> dict:
    """Index-derived parameters — deterministic, no randomness anywhere."""
    return {
        "topic": TOPICS[i % len(TOPICS)],
        "addr": ADDRESSES[i % len(ADDRESSES)],
        "n1": 1 + (i % 4),
        "days": 5 + (i % 9) * 5,
        "days2": 3 + (i % 5),
        "days3": 5 + (i % 6),
        "sum": f"{12_480 + i * 1_137:,}".replace(",", " ") + ",00",
        "vat": 20,
        "pct": 20 + (i % 4) * 10,
        "rate": f"0,{1 + (i % 5)}",
        "cap": 10 + (i % 3) * 5,
        "years": 3 + (i % 3),
        "end_date": f"31 декабря {2024 + (i % 3)} года",
    }


# --------------------------------------------------------------------------
# DOCX builders
# --------------------------------------------------------------------------


#: A real services contract of this length carries a line-item estimate and is
#: printed on company letterhead with a scanned stamp by the signatures. Both
#: are reproduced here: they make the fixture realistic, give the reader view a
#: long table to scroll, and exercise an inline image inside a text flow (a
#: case neither contract-short nor the image-only contract-heavy covers).
ESTIMATE_UNITS = ("шт.", "усл. ед.", "чел.-ч", "компл.", "мес.")
ESTIMATE_STAGES = (
    "обследование",
    "настройка",
    "доработка",
    "тестирование",
    "внедрение",
    "обучение",
    "сопровождение",
    "аудит",
)
ESTIMATE_OBJECTS = (
    "рабочих мест",
    "учётных записей",
    "регламентов",
    "отчётных форм",
    "справочников",
    "интеграций",
    "шаблонов договоров",
    "маршрутов согласования",
)


def _add_estimate_table(doc, rows: int):
    """Смета: one row per work item, every row textually distinct."""
    table = doc.add_table(rows=rows + 1, cols=6)
    table.style = "Table Grid"
    headers = (
        "№",
        "Наименование работ",
        "Ед. изм.",
        "Кол-во",
        "Цена, BYN",
        "Стоимость, BYN",
    )
    for col, text in enumerate(headers):
        cell = table.rows[0].cells[col]
        cell.text = text
        for run in cell.paragraphs[0].runs:
            run.bold = True
    total = 0
    for i in range(rows):
        stage = ESTIMATE_STAGES[i % len(ESTIMATE_STAGES)].capitalize()
        obj = ESTIMATE_OBJECTS[(i * 3) % len(ESTIMATE_OBJECTS)]
        qty = 1 + (i * 7) % 40
        price = 118 + (i * 37) % 2100
        total += qty * price
        values = (
            str(i + 1),
            f"{stage} {obj} (позиция {i + 1:03d}, участок "
            f"{chr(0x410 + i % 32)}-{100 + i})",
            ESTIMATE_UNITS[i % len(ESTIMATE_UNITS)],
            str(qty),
            f"{price},00",
            f"{qty * price},00",
        )
        for col, value in enumerate(values):
            table.rows[i + 1].cells[col].text = value
    return total


def render_letterhead(font_path: str) -> bytes:
    """Company letterhead mark: a flat daisy («ромашка») plus the wordmark."""
    from PIL import Image, ImageDraw, ImageFont

    width, height = 1100, 190
    img = Image.new("RGB", (width, height), (255, 255, 255))
    draw = ImageDraw.Draw(img)
    cx, cy = 95, 95
    for i in range(8):
        angle = i * 45
        petal = Image.new("RGBA", (200, 200), (0, 0, 0, 0))
        ImageDraw.Draw(petal).ellipse((88, 18, 112, 92), fill=(15, 106, 136, 255))
        img.paste(petal.rotate(angle, resample=Image.BICUBIC),
                  (cx - 100, cy - 100), petal.rotate(angle, resample=Image.BICUBIC))
    draw.ellipse((cx - 22, cy - 22, cx + 22, cy + 22), fill=(232, 168, 40))
    draw.text((200, 46), "ООО «Ромашка»", font=ImageFont.truetype(font_path, 46),
              fill=(15, 106, 136))
    draw.text((202, 108),
              "220000, г. Минск, ул. Придуманная, 12 · contracts@example.com",
              font=ImageFont.truetype(font_path, 22), fill=(90, 90, 90))
    draw.line((200, 100, width - 60, 100), fill=(15, 106, 136), width=2)
    out = io.BytesIO()
    img.save(out, format="PNG", optimize=True)
    return out.getvalue()


def render_stamp(font_path: str) -> bytes:
    """A scanned round stamp, the kind pasted next to a signature."""
    from PIL import Image, ImageDraw, ImageFont

    side = 560
    img = Image.new("L", (side, side), 250)
    draw = ImageDraw.Draw(img)
    draw.ellipse((28, 28, side - 28, side - 28), outline=70, width=7)
    draw.ellipse((66, 66, side - 66, side - 66), outline=70, width=3)
    font_big = ImageFont.truetype(font_path, 44)
    font_small = ImageFont.truetype(font_path, 30)
    for i, line in enumerate(("ООО", "«РОМАШКА»")):
        draw.text((side / 2, 190 + i * 54), line, font=font_big, fill=60,
                  anchor="mm")
    for i, line in enumerate(("УНП 100000001", "г. Минск", "ДЛЯ ДОГОВОРОВ")):
        draw.text((side / 2, 320 + i * 42), line, font=font_small, fill=75,
                  anchor="mm")
    noise = Image.frombytes("L", (side, side), keystream("stamp", side * side))
    img = Image.blend(img, noise, 0.08)
    img = img.rotate(-4.5, resample=Image.BICUBIC, fillcolor=250)
    out = io.BytesIO()
    img.save(out, format="JPEG", quality=76, optimize=True)
    return out.getvalue()


def _docx_new(title: str):
    from docx import Document

    doc = Document()
    props = doc.core_properties
    props.title = title
    props.author = "Отдел договорной работы"
    props.last_modified_by = "Отдел договорной работы"
    props.comments = "Тестовая фикстура. Все данные вымышлены."
    props.created = FIXED_DT
    props.modified = FIXED_DT
    props.revision = 1
    return doc


def _add_services_table(doc, rows: int, start: int = 0):
    """Small price table: № / услуга / кол-во / цена / стоимость."""
    table = doc.add_table(rows=rows + 1, cols=5)
    table.style = "Table Grid"
    headers = ("№", "Наименование услуги", "Кол-во", "Цена, BYN", "Стоимость, BYN")
    for col, text in enumerate(headers):
        cell = table.rows[0].cells[col]
        cell.text = text
        for run in cell.paragraphs[0].runs:
            run.bold = True
    for r in range(rows):
        i = start + r
        qty = 1 + (i % 5)
        price = 340 + (i % 7) * 115
        cells = table.rows[r + 1].cells
        cells[0].text = str(r + 1)
        cells[1].text = SERVICE_NAMES[i % len(SERVICE_NAMES)]
        cells[2].text = str(qty)
        cells[3].text = f"{price},00"
        cells[4].text = f"{qty * price},00"
    return table


def _add_signatures(doc):
    doc.add_heading("Реквизиты и подписи Сторон", level=1)
    table = doc.add_table(rows=len(REQUISITES) + 1, cols=2)
    table.style = "Table Grid"
    table.rows[0].cells[0].text = "Исполнитель"
    table.rows[0].cells[1].text = "Заказчик"
    for r, (label, value) in enumerate(REQUISITES, start=1):
        table.rows[r].cells[0].text = f"{label}: {value}"
        table.rows[r].cells[1].text = f"{label}: {CUSTOMER_REQUISITES[label]}"
    doc.add_paragraph()
    doc.add_paragraph("Директор ____________________ / П. П. Петров /")
    doc.add_paragraph("Заказчик ____________________ / И. И. Иванов /")


def _contract_preamble(doc, number: str, title: str):
    heading = doc.add_heading(title, level=0)
    heading.alignment = 1  # WD_ALIGN_PARAGRAPH.CENTER
    subtitle = doc.add_paragraph(f"№ {number} от 15 января 2024 г.")
    subtitle.alignment = 1
    doc.add_paragraph("г. Минск")
    doc.add_paragraph(
        f"{EXECUTOR}, именуемое в дальнейшем «Исполнитель», в лице директора "
        "Петрова Петра Петровича, действующего на основании Устава, с одной "
        f"стороны, и {CUSTOMER}, именуемый в дальнейшем «Заказчик», "
        "действующий на основании свидетельства о государственной регистрации, "
        "с другой стороны, совместно именуемые «Стороны», заключили настоящий "
        "Договор о нижеследующем."
    )


def _add_sections(doc, section_titles, start_no: int = 1, repeats: int = 1):
    """Write numbered sections.

    Clause numbers are written into the text («2.3. …») instead of using Word's
    list numbering: real Russian contracts are typed that way, and it means the
    numbers survive converters that drop numbering definitions (mammoth does).

    ``repeats`` cycles the clause pool with fresh parameters, which is how the
    long contract reaches its page count without a wall of identical sentences.
    """
    number = start_no
    counter = 0
    for title in section_titles:
        clauses = CLAUSES[title]
        doc.add_heading(f"{number}. {title}", level=1)
        sub = 1
        for _rep in range(repeats):
            for clause in clauses:
                doc.add_paragraph(
                    f"{number}.{sub}. " + clause.format(**clause_params(counter))
                )
                sub += 1
                counter += 1
        number += 1
    return number


def build_contract_short(path: Path) -> None:
    doc = _docx_new("Договор возмездного оказания услуг № 17/2024")
    _contract_preamble(doc, "17/2024", "ДОГОВОР ВОЗМЕЗДНОГО ОКАЗАНИЯ УСЛУГ")
    _add_sections(
        doc,
        [
            "Предмет договора",
            "Стоимость услуг и порядок расчётов",
            "Ответственность Сторон",
            "Прочие условия",
        ],
    )
    doc.add_heading("Приложение № 1. Перечень услуг", level=1)
    _add_services_table(doc, rows=4)
    doc.add_paragraph()
    _add_signatures(doc)
    doc.save(path)
    normalize_zip(path, slim_docx)


def build_contract_long(path: Path, letterhead: bytes = None,
                        stamp: bytes = None) -> None:
    from docx.shared import Cm

    doc = _docx_new("Договор на оказание услуг № 44/2024 (расширенная редакция)")
    if letterhead is not None:
        doc.add_picture(io.BytesIO(letterhead), width=Cm(16.0))
    _contract_preamble(
        doc, "44/2024", "ДОГОВОР НА ОКАЗАНИЕ УСЛУГ (РАСШИРЕННАЯ РЕДАКЦИЯ)"
    )
    next_no = _add_sections(
        doc, list(CLAUSES.keys()), start_no=1, repeats=LONG_CLAUSE_REPEATS
    )

    # Appendices: each starts on a new page and carries its own headings,
    # prose and a table, so the reader view has plenty to scroll through.
    subheads = (
        "Состав работ",
        "Требования к результату",
        "Порядок взаимодействия",
        "Ограничения и допущения",
    )
    for appendix in range(1, LONG_APPENDIX_COUNT + 1):
        doc.add_page_break()
        doc.add_heading(f"Приложение № {appendix} к Договору № 44/2024", level=1)
        doc.add_paragraph(
            "Настоящее Приложение определяет объём услуг по этапу "
            f"№ {appendix} и подписано уполномоченными представителями Сторон "
            "15 января 2024 г. Ответственный исполнитель со стороны "
            f"{EXECUTOR} — ведущий специалист "
            f"{'Сидоров С. С.' if appendix % 2 else 'Кузнецова А. В.'}."
        )
        for part in range(1, 5):
            doc.add_heading(f"{appendix}.{part}. {subheads[part - 1]}", level=2)
            for k in range(6):
                idx = appendix * 37 + part * 11 + k
                params = clause_params(idx)
                doc.add_paragraph(
                    f"{appendix}.{part}.{k + 1}. Исполнитель выполняет работы по "
                    f"направлению «{TOPICS[idx % len(TOPICS)]}» в срок "
                    f"{params['days']} календарных дней, привлекая при "
                    f"необходимости {THIRD_PARTIES[idx % len(THIRD_PARTIES)]} на "
                    "условиях, согласованных с Заказчиком. Стоимость этапа "
                    f"составляет {params['sum']} белорусских рублей, включая "
                    f"НДС {params['vat']}%. Отчётные документы предоставляются "
                    f"по адресу {params['addr']} не позднее "
                    f"{params['days3']} рабочих дней с даты завершения этапа."
                )
        _add_services_table(doc, rows=6, start=appendix * 3)
        doc.add_paragraph()

    # Смета: a long line-item table, the part a reader scrolls forever.
    doc.add_page_break()
    doc.add_heading(
        f"Приложение № {LONG_APPENDIX_COUNT + 1}. Смета на выполнение работ",
        level=1,
    )
    doc.add_paragraph(
        "Смета составлена в белорусских рублях и включает налог на добавленную "
        "стоимость по ставке 20%. Позиции сметы соответствуют этапам, "
        "перечисленным в приложениях к настоящему Договору."
    )
    total = _add_estimate_table(doc, LONG_ESTIMATE_ROWS)
    doc.add_paragraph(
        f"Итого по смете: {total:,}".replace(",", " ") + ",00 белорусских рублей."
    )

    # Extra sections after the appendices, continuing the numbering.
    _add_sections(
        doc,
        [
            "Конфиденциальность",
            "Обработка персональных данных",
            "Порядок разрешения споров",
            "Прочие условия",
        ],
        start_no=next_no,
        repeats=2,
    )
    doc.add_page_break()
    _add_signatures(doc)
    if stamp is not None:
        doc.add_paragraph("М. П.")
        doc.add_picture(io.BytesIO(stamp), width=Cm(4.2))
    doc.save(path)
    normalize_zip(path, slim_docx)


def build_contract_heavy(path: Path, images: list) -> int:
    """Contract with scanned appendices until the file clears the CDN cutoff.

    Text compresses far too well to reach 350 KB on its own, so each appendix
    embeds an already-compressed JPEG "scan" the zip cannot squeeze further.
    The loop is deterministic: same inputs, same number of pages.
    """
    from docx.shared import Cm

    size = 0
    for count in range(1, len(images) + 1):
        doc = _docx_new("Договор подряда № 77/2024 со сканированными приложениями")
        _contract_preamble(doc, "77/2024", "ДОГОВОР ПОДРЯДА")
        _add_sections(
            doc,
            [
                "Предмет договора",
                "Права и обязанности Сторон",
                "Стоимость услуг и порядок расчётов",
                "Порядок сдачи и приёмки услуг",
            ],
        )
        for i in range(count):
            doc.add_page_break()
            doc.add_heading(
                f"Приложение № {i + 1}. Скан-копия акта (лист {i + 1})", level=1
            )
            doc.add_paragraph(
                "Скан-копия приведена для сведения. Оригинал документа хранится "
                f"у Исполнителя по адресу {ADDRESSES[i % len(ADDRESSES)]}."
            )
            doc.add_picture(io.BytesIO(images[i]), width=Cm(15.5))
        _add_signatures(doc)
        doc.save(path)
        normalize_zip(path, slim_docx)
        size = path.stat().st_size
        if size >= HEAVY_TARGET_BYTES:
            break
    return size


def build_broken_docx(path: Path) -> None:
    """A truncated docx: right magic bytes, unreadable archive.

    Keeps ``PK\\x03\\x04`` so extension-and-sniffing checks still call it a
    docx, but the central directory is gone, so every zip reader fails.
    """
    doc = _docx_new("Договор № 99/2024 (повреждённый файл)")
    _contract_preamble(doc, "99/2024", "ДОГОВОР ОКАЗАНИЯ УСЛУГ")
    _add_sections(doc, ["Предмет договора", "Прочие условия"])
    intact = path.with_name(path.name + ".intact")
    doc.save(intact)
    normalize_zip(intact, slim_docx)  # same fixed timestamps as every other docx
    whole = intact.read_bytes()
    intact.unlink()
    # Cut the archive off mid-stream: the local file headers survive (so the
    # magic bytes and the first entries still look like a docx) but the central
    # directory is gone, which is what makes every zip reader refuse it.
    path.write_bytes(whole[: int(len(whole) * 0.55)])


def build_locked_docx(path: Path) -> str:
    """Password-protected docx via msoffcrypto-tool (optional dependency).

    The output is a real CFB/OLE2 container with ECMA-376 agile encryption —
    not a renamed zip. Returns a status string for the report.
    """
    try:
        import msoffcrypto
        import msoffcrypto.format.ooxml
    except ImportError:
        return "SKIPPED (msoffcrypto-tool not installed)"

    if not hasattr(msoffcrypto.format.ooxml.OOXMLFile, "encrypt"):
        return "SKIPPED (installed msoffcrypto-tool can only decrypt, not encrypt)"

    doc = _docx_new("Договор № 55/2024 (защищён паролем)")
    _contract_preamble(doc, "55/2024", "ДОГОВОР ОКАЗАНИЯ УСЛУГ")
    _add_sections(doc, ["Предмет договора", "Конфиденциальность"])
    doc.add_paragraph("Документ защищён паролем. Пароль для тестов: " + DOCX_PASSWORD)
    plain = io.BytesIO()
    doc.save(plain)
    plain.seek(0)
    with open(path, "wb") as out:
        msoffcrypto.OfficeFile(plain).encrypt(DOCX_PASSWORD, out)
    return "OK"


# --------------------------------------------------------------------------
# XLSX builders
# --------------------------------------------------------------------------

CONTRACTORS = (
    "ООО «Ромашка»",
    "ИП Иванов И. И.",
    "ООО «Василёк»",
    "ЗАО «Тестовый Мост»",
    "ООО «ТехноСфера-Плюс»",
    "ОДО «Северный Ветер»",
    "ООО «Гранит-Строй»",
)
STATUSES = ("Действует", "На согласовании", "Завершён", "Продлён", "Расторгнут")
SUBJECTS = (
    "Оказание услуг",
    "Поставка оборудования",
    "Техническое обслуживание",
    "Аренда помещения",
    "Консультационные услуги",
    "Подряд на монтаж",
)


def _xlsx_props(wb, title: str):
    wb.properties.title = title
    wb.properties.creator = "Отдел договорной работы"
    wb.properties.lastModifiedBy = "Отдел договорной работы"
    wb.properties.created = FIXED_DT
    wb.properties.modified = FIXED_DT


def cache_formula_results(names, data):
    """Give some of sheet 2's formulas a cached result, and leave others without.

    openpyxl writes ``<f>C4-B4</f><v></v>`` — an *empty* cached value, because
    it never evaluates formulas. Viewers do not evaluate them either: SheetJS (the
    library behind the spreadsheet preview) drops a formula cell that has no
    cached value, so such a cell renders as an empty td. Excel, by contrast,
    always stores the last computed result.

    Both shapes occur in files users upload, so sheet 2 carries both on
    purpose: column D ("Отклонение") and the totals row get cached values and
    render as numbers; column E ("Признак") stays uncached and renders empty.
    A preview test can assert exactly that difference.
    """
    part = "xl/worksheets/sheet2.xml"
    if part not in data:
        return names, data
    cached = {
        "D4": "-1569.5",
        "D5": "3120",
        "D6": "-3995.75",
        "D7": "-119.25",
        "B8": "277000",
        "C8": "274435.5",
        "D8": "-2564.5",
    }
    xml = data[part].decode("utf-8")
    for ref, value in cached.items():
        xml = re.sub(
            rf'(<c r="{ref}"[^>]*>)(<f>[^<]*</f>)(?:<v>\s*</v>)?(</c>)',
            rf"\g<1>\g<2><v>{value}</v>\g<3>",
            xml,
        )
    data[part] = xml.encode("utf-8")
    return names, data


def build_registry_xlsx(path: Path) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    _xlsx_props(wb, "Реестр договоров на 2024 год")

    ws = wb.active
    ws.title = "Реестр договоров"

    # A1 title + a deliberately blank row 2 before the header block.
    ws["A1"] = "Реестр действующих договоров ООО «Ромашка» на 2024 год"
    ws["A1"].font = Font(bold=True, size=14)
    ws.merge_cells("A1:H1")
    ws["A1"].alignment = Alignment(horizontal="center")

    # Two-level header with merged groups (rows 3-4).
    ws["A3"] = "№"
    ws.merge_cells("A3:A4")
    ws["B3"] = "Контрагент"
    ws.merge_cells("B3:B4")
    ws["C3"] = "Предмет договора"
    ws.merge_cells("C3:C4")
    ws["D3"] = "Сроки"
    ws.merge_cells("D3:E3")
    ws["D4"] = "Начало"
    ws["E4"] = "Окончание"
    ws["F3"] = "Суммы, BYN"
    ws.merge_cells("F3:G3")
    ws["F4"] = "Сумма договора"
    ws["G4"] = "Оплачено"
    ws["H3"] = "Статус"
    ws.merge_cells("H3:H4")
    for row in (3, 4):
        for col in range(1, 9):
            cell = ws.cell(row=row, column=col)
            cell.font = Font(bold=True)
            cell.alignment = Alignment(
                horizontal="center", vertical="center", wrap_text=True
            )

    total_amount = 0.0
    total_paid = 0.0
    first_data_row = 5
    rows = 40
    for i in range(rows):
        r = first_data_row + i
        amount = 4_800 + i * 1_130 + (i % 7) * 215
        paid = round(amount * (0.25 + (i % 4) * 0.25), 2)
        ws.cell(row=r, column=1, value=i + 1)
        ws.cell(row=r, column=2, value=CONTRACTORS[i % len(CONTRACTORS)])
        ws.cell(row=r, column=3, value=SUBJECTS[i % len(SUBJECTS)])
        ws.cell(row=r, column=4, value=f"{1 + i % 28:02d}.0{1 + i % 9}.2024")
        # Deliberate holes: every 5th row has no end date, every 7th no status,
        # every 11th no paid amount. The preview must not shift columns.
        if i % 5 != 0:
            ws.cell(row=r, column=5, value=f"{1 + (i * 3) % 28:02d}.12.2024")
        ws.cell(row=r, column=6, value=float(amount))
        if i % 11 != 0:
            ws.cell(row=r, column=7, value=float(paid))
            total_paid += paid
        if i % 7 != 0:
            ws.cell(row=r, column=8, value=STATUSES[i % len(STATUSES)])
        total_amount += amount

    total_row = first_data_row + rows
    ws.cell(row=total_row, column=1, value="Итого")
    ws.merge_cells(
        start_row=total_row, start_column=1, end_row=total_row, end_column=5
    )
    ws.cell(row=total_row, column=1).font = Font(bold=True)
    ws.cell(row=total_row, column=1).alignment = Alignment(horizontal="right")
    # Literal numbers, not formulas: a viewer that does not recalculate still
    # shows the totals (see sheet 2 for the deliberate formula case).
    ws.cell(row=total_row, column=6, value=round(total_amount, 2)).font = Font(bold=True)
    ws.cell(row=total_row, column=7, value=round(total_paid, 2)).font = Font(bold=True)
    ws.cell(row=total_row, column=8, value="—")

    for idx, width in enumerate((6, 30, 26, 12, 12, 16, 14, 18), start=1):
        ws.column_dimensions[get_column_letter(idx)].width = width
    ws.freeze_panes = "A5"

    # Sheet 2 — formulas. openpyxl writes no cached values, so a viewer that
    # does not evaluate formulas renders these cells empty. That is on purpose:
    # it is the "spreadsheet with formulas" case the preview has to survive.
    ws2 = wb.create_sheet("Сводка")
    ws2["A1"] = "Сводка по кварталам"
    ws2["A1"].font = Font(bold=True, size=12)
    for cell, text in (
        ("A3", "Квартал"),
        ("B3", "План, BYN"),
        ("C3", "Факт, BYN"),
        ("D3", "Отклонение"),
        ("E3", "Признак"),
    ):
        ws2[cell] = text
        ws2[cell].font = Font(bold=True)
    plan_fact = (
        (60_000, 58_430.5),
        (72_000, 75_120.0),
        (65_000, 61_004.25),
        (80_000, 79_880.75),
    )
    for i, (plan, fact) in enumerate(plan_fact):
        r = 4 + i
        ws2.cell(row=r, column=1, value=f"{i + 1} квартал")
        ws2.cell(row=r, column=2, value=float(plan))
        ws2.cell(row=r, column=3, value=float(fact))
        ws2.cell(row=r, column=4, value=f"=C{r}-B{r}")
        ws2.cell(row=r, column=5, value=f'=IF(C{r}>=B{r},"выполнен","не выполнен")')
    ws2["A8"] = "Итого"
    ws2["A8"].font = Font(bold=True)
    ws2["B8"] = "=SUM(B4:B7)"
    ws2["C8"] = "=SUM(C4:C7)"
    ws2["D8"] = "=SUM(D4:D7)"
    for col, width in zip("ABCDE", (14, 14, 14, 14, 16)):
        ws2.column_dimensions[col].width = width

    # Sheet 3 — nearly empty.
    ws3 = wb.create_sheet("Черновик")
    ws3["B2"] = "Лист намеренно оставлен почти пустым."

    wb.save(path)
    normalize_zip(path, cache_formula_results)


def build_big_rows_xlsx(path: Path, data_rows: int = 6_000) -> None:
    """One sheet well past SPREADSHEET_MAX_ROWS_PER_SHEET (5 000)."""
    from openpyxl import Workbook
    from openpyxl.styles import Font
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    _xlsx_props(wb, "Журнал операций за 2024 год")
    ws = wb.active
    ws.title = "Журнал операций"

    headers = ("№", "Дата", "Контрагент", "Документ", "Сумма, BYN", "Статус")
    ws.append(headers)
    for col in range(1, len(headers) + 1):
        ws.cell(row=1, column=col).font = Font(bold=True)

    for i in range(data_rows):
        ws.append(
            (
                i + 1,
                f"{1 + i % 28:02d}.{1 + i % 12:02d}.2024",
                CONTRACTORS[i % len(CONTRACTORS)],
                f"АКТ-{100000 + i}",
                float(120 + (i % 997) * 13),
                STATUSES[i % len(STATUSES)],
            )
        )
    for idx, width in enumerate((8, 12, 30, 16, 14, 18), start=1):
        ws.column_dimensions[get_column_letter(idx)].width = width
    ws.freeze_panes = "A2"
    wb.save(path)
    normalize_zip(path)


# --------------------------------------------------------------------------
# PPTX builders
# --------------------------------------------------------------------------

DECK_SLIDES = (
    (
        "Итоги квартала",
        (
            "Выручка выросла на 12% к предыдущему кварталу",
            "Подписано 8 новых договоров",
            "Средний срок согласования — 6 рабочих дней",
        ),
    ),
    ("Структура портфеля", ("Услуги — 54%", "Поставка — 31%", "Аренда — 15%")),
    (
        "Ключевые контрагенты",
        ("ООО «Василёк»", "ЗАО «Тестовый Мост»", "ООО «ТехноСфера-Плюс»"),
    ),
    (
        "Просроченная задолженность",
        (
            "Всего 4 договора",
            "Средний срок просрочки — 11 дней",
            "Претензионная работа начата по 3 из 4",
        ),
    ),
    (
        "Процесс согласования",
        ("Заявка", "Проверка контрагента", "Юридическая экспертиза", "Подписание"),
    ),
    (
        "Риски",
        (
            "Зависимость от одного поставщика",
            "Валютные колебания",
            "Сроки поставки оборудования",
        ),
    ),
    (
        "План на следующий квартал",
        (
            "Перевести 60% договоров в электронный вид",
            "Сократить срок согласования до 4 дней",
            "Запустить реестр в общем доступе",
        ),
    ),
    (
        "Бюджет проекта",
        (
            "Расходы на лицензии — 18 400 BYN",
            "Обучение — 6 200 BYN",
            "Резерв — 5 000 BYN",
        ),
    ),
    (
        "Команда",
        (
            "Руководитель проекта — П. П. Петров",
            "Аналитик — А. В. Кузнецова",
            "Юрист — С. С. Сидоров",
        ),
    ),
    (
        "Метрики (KPI)",
        ("Time to sign — 6 дней", "Contract coverage — 92%", "Renewal rate — 78%"),
    ),
)


def _pptx_new(width_emu: int, height_emu: int, title: str):
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width = Emu(width_emu)
    prs.slide_height = Emu(height_emu)
    props = prs.core_properties
    props.title = title
    props.author = "Отдел договорной работы"
    props.last_modified_by = "Отдел договорной работы"
    props.created = FIXED_DT
    props.modified = FIXED_DT
    props.revision = 1
    return prs


def _add_title_slide(prs, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def _add_bullets_slide(prs, title: str, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.text = bullets[0]
    for line in bullets[1:]:
        para = body.add_paragraph()
        para.text = line
        para.level = 1
    return slide


def _add_table_slide(prs, title: str):
    from pptx.util import Cm, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    rows, cols = 5, 4
    table = slide.shapes.add_table(
        rows, cols, Cm(2), Cm(4), prs.slide_width - Cm(4), Cm(8)
    ).table
    for c, text in enumerate(("Договор", "Контрагент", "Сумма, BYN", "Статус")):
        table.cell(0, c).text = text
    for r in range(1, rows):
        i = r - 1
        table.cell(r, 0).text = f"№ {17 + i}/2024"
        table.cell(r, 1).text = CONTRACTORS[i % len(CONTRACTORS)]
        table.cell(r, 2).text = f"{12_480 + i * 3_150:,}".replace(",", " ") + ",00"
        table.cell(r, 3).text = STATUSES[i % len(STATUSES)]
    for row in table.rows:
        for cell in row.cells:
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(14)
    return slide


def _add_barchart_slide(prs, title: str):
    """Bars drawn as shapes — a chart-like slide without an embedded chart part."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Cm, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    values = (42, 68, 55, 91, 74)
    labels = ("Янв", "Фев", "Мар", "Апр", "Май")
    base_top = Cm(13)
    for i, (value, label) in enumerate(zip(values, labels)):
        height = Cm(value / 10)
        left = Cm(3 + i * 3.2)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, base_top - height, Cm(2.2), height
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(0x0F, 0x6A, 0x88)
        bar.line.fill.background()
        caption = slide.shapes.add_textbox(left, base_top, Cm(2.2), Cm(1))
        para = caption.text_frame.paragraphs[0]
        para.text = f"{label} — {value}"
        para.runs[0].font.size = Pt(11)
    return slide


def build_deck_16x9(path: Path) -> None:
    prs = _pptx_new(12_192_000, 6_858_000, "Отчёт по договорной работе за квартал")
    _add_title_slide(
        prs,
        "Отчёт по договорной работе",
        "ООО «Ромашка» · IV квартал 2024 · для внутреннего использования",
    )
    for title, bullets in DECK_SLIDES[:9]:
        _add_bullets_slide(prs, title, bullets)
    _add_table_slide(prs, "Реестр крупных договоров")
    _add_barchart_slide(prs, "Динамика подписаний (шт.)")
    prs.save(path)
    normalize_zip(path)


def build_deck_4x3(path: Path) -> None:
    prs = _pptx_new(9_144_000, 6_858_000, "Краткая сводка 4:3")
    _add_title_slide(prs, "Краткая сводка", "Формат 4:3 · ООО «Ромашка» · январь 2024")
    for title, bullets in DECK_SLIDES[:4]:
        _add_bullets_slide(prs, title, bullets)
    _add_table_slide(prs, "Договоры на контроле")
    prs.save(path)
    normalize_zip(path)


def build_deck_many(path: Path, slides: int = 60) -> None:
    """60 slides, every title unique, so navigation tests have something to
    assert on ("go to slide 47" must land on a distinguishable slide)."""
    prs = _pptx_new(12_192_000, 6_858_000, "Годовой обзор — расширенная версия")
    _add_title_slide(
        prs, "Годовой обзор", "Расширенная версия · 60 слайдов · ООО «Ромашка»"
    )
    for i in range(2, slides + 1):
        base_title, bullets = DECK_SLIDES[(i - 2) % len(DECK_SLIDES)]
        extra = (
            f"Раздел {i:02d}: показатель за период — {1200 + i * 37} условных единиц"
        )
        _add_bullets_slide(
            prs, f"Слайд {i:02d}. {base_title}", tuple(bullets) + (extra,)
        )
    prs.save(path)
    normalize_zip(path)


# --------------------------------------------------------------------------
# Plain-text fixtures
# --------------------------------------------------------------------------

NOTES_MD = """\
# Заметки по договорной работе

Короткий рабочий файл: как мы ведём реестр договоров и что проверяем перед
подписанием. Все примеры вымышленные.

## Что проверяем перед подписанием

1. Полномочия подписанта (устав, доверенность).
2. Реквизиты контрагента — сверяем с карточкой в реестре.
3. Сроки оплаты и порядок приёмки.
4. Ответственность: пеня, предельный размер убытков.

- [x] Шаблон договора обновлён
- [x] Реестр выгружается в XLSX
- [ ] Автонапоминание о продлении

## Реестр: поля

| Поле        | Тип     | Обязательное | Комментарий                     |
|-------------|---------|--------------|---------------------------------|
| Номер       | строка  | да           | Формат `NN/ГГГГ`                |
| Контрагент  | строка  | да           | Например, ООО «Ромашка»         |
| Сумма       | число   | да           | В белорусских рублях, с НДС     |
| Окончание   | дата    | нет          | Пусто — договор бессрочный      |

## Как выгрузить реестр

```python
from registry import export_contracts

# Выгружаем только действующие договоры за 2024 год.
rows = export_contracts(year=2024, status="Действует")
print(f"Выгружено строк: {len(rows)}")
```

> Важно: выгрузка не заменяет сверку взаиморасчётов.
> Сверку проводим ежеквартально и оформляем актом.

Подробности — во внутренней инструкции:
[Регламент договорной работы](https://example.com/regulations/contracts).

---

*Файл-фикстура для тестов предпросмотра. Any resemblance to real companies is
purely coincidental.*
"""

SCRIPT_PY = '''\
#!/usr/bin/env python3
"""Мини-утилита для сводки по реестру договоров.

Фикстура для тестов предпросмотра кода: файл должен читаться как обычный
рабочий скрипт — с докстрингами, комментариями и осмысленными именами.
Данные вымышленные.
"""

import csv
from collections import defaultdict
from decimal import Decimal

# Договор считается крупным, если его сумма не меньше этого порога (BYN).
LARGE_CONTRACT_THRESHOLD = Decimal("50000.00")

# Статусы, которые не учитываем в сводке: договор не порождает обязательств.
IGNORED_STATUSES = frozenset({"Расторгнут", "Черновик"})


def read_registry(path):
    """Читает CSV-реестр и возвращает список словарей.

    Пустая ячейка суммы означает «сумма не согласована» — такие строки
    оставляем, но в деньги не превращаем.
    """
    with open(path, encoding="utf-8", newline="") as handle:
        return list(csv.DictReader(handle))


def parse_amount(raw):
    """Превращает «12 480,00» в Decimal. Пустая строка -> None."""
    if not raw or not raw.strip():
        return None
    normalized = raw.replace("\\u00a0", "").replace(" ", "").replace(",", ".")
    return Decimal(normalized)


def summarize(rows):
    """Считает количество и сумму договоров по каждому контрагенту."""
    totals = defaultdict(lambda: {"count": 0, "amount": Decimal("0")})
    for row in rows:
        if row.get("Статус") in IGNORED_STATUSES:
            continue
        amount = parse_amount(row.get("Сумма"))
        bucket = totals[row["Контрагент"]]
        bucket["count"] += 1
        if amount is not None:
            bucket["amount"] += amount
    return totals


def large_contracts(rows):
    """Возвращает договоры не меньше порога, от большего к меньшему."""
    result = []
    for row in rows:
        amount = parse_amount(row.get("Сумма"))
        if amount is not None and amount >= LARGE_CONTRACT_THRESHOLD:
            result.append((amount, row["Номер"], row["Контрагент"]))
    return sorted(result, reverse=True)


def main(path="data.csv"):
    """Печатает сводку по контрагентам и список крупных договоров."""
    rows = read_registry(path)
    totals = summarize(rows)

    print(f"Контрагентов в реестре: {len(totals)}")
    for name, bucket in sorted(totals.items()):
        print(f"  {name:<40} {bucket['count']:>3} шт.  {bucket['amount']:>12}")

    print("\\nКрупные договоры:")
    for amount, number, name in large_contracts(rows):
        print(f"  {number:<10} {name:<40} {amount:>12}")


if __name__ == "__main__":
    main()
'''

CSV_SUBJECTS = (
    "Оказание услуг, включая консультации",  # comma inside the field
    'Поставка оборудования "Тестовый Мост"',  # quotes inside the field
    "Аренда помещения",
    "Техническое обслуживание, 12 месяцев",
    "Монтаж и пусконаладка",
    'Услуги по договору "под ключ", этап 1',  # both comma and quotes
)


def build_data_csv(path: Path, rows: int = 30) -> None:
    header = ("Номер", "Контрагент", "Предмет", "Сумма", "Статус", "Комментарий")
    with open(path, "w", encoding="utf-8", newline="") as handle:
        writer = csv.writer(handle, lineterminator="\n")
        writer.writerow(header)
        for i in range(rows):
            amount = 4_800 + i * 1_970
            # Every 4th comment is empty on purpose (empty field case).
            comment = "" if i % 4 == 0 else f"Проверено {1 + i % 28:02d}.01.2024"
            writer.writerow(
                (
                    f"{17 + i}/2024",
                    CONTRACTORS[i % len(CONTRACTORS)],
                    CSV_SUBJECTS[i % len(CSV_SUBJECTS)],
                    f"{amount},00",
                    STATUSES[i % len(STATUSES)],
                    comment,
                )
            )


def build_archive_zip(path: Path) -> None:
    """Valid zip with two text files — an unsupported type for the preview."""
    members = {
        "readme.txt": (
            "Архив-фикстура для тестов предпросмотра.\n"
            "Предпросмотр архивов не поддерживается — ожидается заглушка "
            "с предложением скачать файл.\n"
        ),
        "contracts/список.txt": (
            "Договор № 17/2024 — ООО «Ромашка»\n"
            "Договор № 18/2024 — ИП Иванов И. И.\n"
        ),
    }
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED, compresslevel=9) as zf:
        for name, text in members.items():
            info = zipfile.ZipInfo(name, date_time=ZIP_DT)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = ZIP_ATTR
            info.create_system = 0
            zf.writestr(info, text.encode("utf-8"))


def build_unknown_xyz(path: Path) -> None:
    """Opaque binary blob with an extension nothing claims to handle."""
    header = b"\x00XYZ\x01FIXTURE\x00"
    path.write_bytes(header + keystream("unknown.xyz", 1024 - len(header)))


# --------------------------------------------------------------------------
# PDF builders
# --------------------------------------------------------------------------

PDF_SECTIONS = (
    (
        "Раздел 1. Общие положения",
        (
            "Настоящий документ определяет порядок ведения договорной работы в "
            "ООО «Ромашка» и применяется ко всем подразделениям организации. "
            "Документ подготовлен исключительно для тестов предпросмотра; все "
            "наименования, суммы и адреса вымышлены.",
            "Ответственным за ведение реестра договоров является отдел "
            "договорной работы. Реестр ведётся в электронном виде и "
            "выгружается в формате XLSX не реже одного раза в квартал.",
            "Термины и определения, используемые в настоящем документе, "
            "применяются в значении, установленном законодательством и "
            "внутренними документами организации.",
        ),
    ),
    (
        "Раздел 2. Согласование договоров",
        (
            "Проект договора направляется на согласование не позднее чем за "
            "пять рабочих дней до предполагаемой даты подписания. Срок "
            "согласования для типовых форм — два рабочих дня.",
            "Согласование проводится последовательно: инициатор, юридическая "
            "служба, финансовая служба, руководитель. Каждый согласующий "
            "фиксирует замечания в листе согласования.",
            "Отклонение от типовой формы допускается только с письменного "
            "разрешения руководителя организации и оформляется отдельной "
            "служебной запиской.",
        ),
    ),
    (
        "Раздел 3. Проверка контрагента",
        (
            "Перед подписанием договора проверяются: правоспособность "
            "контрагента, полномочия подписанта, отсутствие сведений о "
            "ликвидации, а также наличие задолженности по ранее заключённым "
            "договорам.",
            "Результат проверки оформляется карточкой контрагента и хранится "
            "вместе с оригиналом договора в течение всего срока его действия и "
            "трёх лет после прекращения.",
            "Contractor due diligence is repeated annually for long-running "
            "agreements — часть проверок выполняется автоматически по данным "
            "реестра.",
        ),
    ),
    (
        "Раздел 4. Исполнение и приёмка",
        (
            "Приёмка оказанных услуг оформляется актом сдачи-приёмки в двух "
            "экземплярах. Мотивированный отказ от приёмки направляется в "
            "течение трёх рабочих дней с даты получения акта.",
            "Оплата производится в течение пяти банковских дней с даты "
            "подписания акта, если договором не установлен иной срок. "
            "Предварительная оплата допускается в размере до тридцати "
            "процентов стоимости услуг.",
            "Сверка взаиморасчётов проводится ежеквартально; акт сверки "
            "подписывается уполномоченными представителями Сторон.",
        ),
    ),
    (
        "Раздел 5. Хранение и архив",
        (
            "Оригиналы договоров хранятся по адресу: 220000, г. Минск, "
            "ул. Придуманная, д. 12, оф. 305. Доступ к архиву предоставляется "
            "по письменной заявке.",
            "Электронные копии хранятся в системе документооборота с "
            "сохранением истории изменений. Срок хранения электронных копий "
            "совпадает со сроком хранения оригиналов.",
            "Уничтожение документов по истечении срока хранения оформляется "
            "актом, который подписывается комиссией из трёх сотрудников.",
        ),
    ),
)


def _wrap(canvas_obj, text: str, font: str, size: float, max_width: float):
    """Greedy word wrap using the font's real metrics."""
    lines, current = [], ""
    for word in text.split():
        candidate = f"{current} {word}".strip()
        if canvas_obj.stringWidth(candidate, font, size) <= max_width:
            current = candidate
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines


def build_digital_pdf(path: Path, regular: str, bold: str) -> None:
    """Five pages with a genuine text layer.

    reportlab embeds a ToUnicode CMap for the subset it writes, so the Cyrillic
    really is extractable by a viewer — that is what makes «Раздел 3»
    searchable in a test rather than just visible.
    """
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.pdfgen import canvas as rl_canvas

    pdfmetrics.registerFont(TTFont("Body", regular))
    pdfmetrics.registerFont(TTFont("BodyBold", bold))

    width, height = A4
    margin = 56
    c = rl_canvas.Canvas(str(path), pagesize=A4, invariant=1)
    c.setTitle("Регламент договорной работы")
    c.setAuthor("Отдел договорной работы")
    c.setSubject("Фикстура для тестов предпросмотра PDF")

    for heading, paragraphs in PDF_SECTIONS:
        y = height - margin
        c.setFont("BodyBold", 16)
        c.drawString(margin, y, heading)
        y -= 30
        c.setFont("Body", 11)
        for paragraph in paragraphs:
            for line in _wrap(c, paragraph, "Body", 11, width - 2 * margin):
                c.drawString(margin, y, line)
                y -= 16
            y -= 8
        c.setFont("Body", 9)
        c.drawString(margin, 40, "ООО «Ромашка» · Регламент договорной работы")
        c.showPage()
    c.save()


def build_locked_pdf(path: Path, regular: str, bold: str) -> None:
    """Two-page PDF, then AES-256 encrypted with a user password."""
    import pikepdf
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.pdfgen import canvas as rl_canvas

    pdfmetrics.registerFont(TTFont("Body", regular))
    pdfmetrics.registerFont(TTFont("BodyBold", bold))

    width, height = A4
    margin = 56
    buf = io.BytesIO()
    c = rl_canvas.Canvas(buf, pagesize=A4, invariant=1)
    c.setTitle("Конфиденциальное приложение")
    pages = (
        (
            "Конфиденциальное приложение № 1",
            (
                "Документ защищён паролем и предназначен для проверки "
                "поведения предпросмотра при запросе пароля.",
                "Пароль для тестов указан в файле README.md рядом с "
                "фикстурами. Содержимое вымышлено.",
            ),
        ),
        (
            "Конфиденциальное приложение № 2",
            (
                "Стоимость услуг по договору № 55/2024 составляет "
                "18 400,00 белорусских рублей, включая НДС 20%.",
                "Ответственный исполнитель — ведущий специалист "
                "С. С. Сидоров, контактный адрес: sidorov@example.com.",
            ),
        ),
    )
    for heading, paragraphs in pages:
        y = height - margin
        c.setFont("BodyBold", 16)
        c.drawString(margin, y, heading)
        y -= 30
        c.setFont("Body", 11)
        for paragraph in paragraphs:
            for line in _wrap(c, paragraph, "Body", 11, width - 2 * margin):
                c.drawString(margin, y, line)
                y -= 16
            y -= 8
        c.showPage()
    c.save()

    buf.seek(0)
    with pikepdf.open(buf) as pdf:
        pdf.save(
            path,
            encryption=pikepdf.Encryption(
                user=PDF_PASSWORD, owner="owner-" + PDF_PASSWORD, R=6
            ),
        )


SCAN_PAGES = (
    (
        "АКТ сдачи-приёмки оказанных услуг № 3",
        "",
        "г. Минск                                  15 января 2024 г.",
        "",
        "ООО «Ромашка», именуемое «Исполнитель», с одной стороны,",
        "и ИП Иванов И. И., именуемый «Заказчик», с другой стороны,",
        "составили настоящий акт о нижеследующем.",
        "",
        "1. Услуги оказаны в полном объёме и в согласованный срок.",
        "2. Стоимость услуг — 12 480,00 BYN, включая НДС 20%.",
        "3. Заказчик претензий по объёму и качеству не имеет.",
    ),
    (
        "Приложение к акту № 3. Перечень работ",
        "",
        "1. Обследование процессов — 24 часа.",
        "2. Разработка регламента — 40 часов.",
        "3. Настройка рабочих мест — 16 часов.",
        "4. Обучение сотрудников — 8 часов.",
        "",
        "Итого: 88 часов.",
        "",
        "Работы выполнены по адресу: 220000, г. Минск,",
        "ул. Придуманная, д. 12, оф. 305.",
    ),
    (
        "Письмо о согласовании № 41/24",
        "",
        "Уважаемые коллеги!",
        "",
        "Подтверждаем согласование редакции договора № 44/2024",
        "в части сроков оплаты и порядка приёмки услуг.",
        "",
        "Просим направить подписанный экземпляр по адресу:",
        "contracts@example.com.",
        "",
        "С уважением, отдел договорной работы ООО «Ромашка».",
    ),
)


def render_scan_page(
    page_no: int, lines, font_path: str, size=(827, 1169), quality: int = 45
) -> bytes:
    """Render a page as a grayscale JPEG that looks like a flatbed scan.

    Deliberately raster-only: dropped into a PDF this produces a document with
    no text layer at all, which is what makes it a useful "scanned document"
    fixture. Slight rotation plus deterministic speckle keeps it from
    compressing away and makes it look scanned rather than generated.
    """
    from PIL import Image, ImageDraw, ImageFont

    page = Image.new("L", size, 247)
    draw = ImageDraw.Draw(page)
    title_font = ImageFont.truetype(font_path, 30)
    body_font = ImageFont.truetype(font_path, 21)

    y = 96
    draw.text((90, y), lines[0], font=title_font, fill=35)
    y += 60
    for line in lines[1:]:
        draw.text((90, y), line, font=body_font, fill=45)
        y += 34
    # A stamp-like box and a signature line, the way scans usually look.
    draw.rectangle(
        (size[0] - 300, size[1] - 320, size[0] - 90, size[1] - 190), outline=90, width=3
    )
    draw.text((size[0] - 285, size[1] - 285), "КОПИЯ ВЕРНА", font=body_font, fill=90)
    draw.text((size[0] - 285, size[1] - 245), f"лист {page_no}", font=body_font, fill=90)
    draw.line((90, size[1] - 150, 420, size[1] - 150), fill=110, width=2)
    draw.text((90, size[1] - 140), "подпись / дата", font=body_font, fill=120)

    # Deterministic speckle: blend in a noise plane at low weight.
    noise = Image.frombytes("L", size, keystream(f"scan-{page_no}", size[0] * size[1]))
    page = Image.blend(page, noise, 0.055)
    # Scans are never perfectly straight.
    page = page.rotate(
        0.6 if page_no % 2 else -0.45, resample=Image.BICUBIC, fillcolor=246
    )

    out = io.BytesIO()
    page.save(out, format="JPEG", quality=quality, optimize=True)
    return out.getvalue()


def build_scan_pdf(path: Path, font_path: str) -> None:
    """Three image-only pages: no fonts, no text operators, nothing to extract."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas as rl_canvas

    width, height = A4
    c = rl_canvas.Canvas(str(path), pagesize=A4, invariant=1)
    c.setTitle("Скан-копия документов")
    for i, lines in enumerate(SCAN_PAGES, start=1):
        jpeg = render_scan_page(i, lines, font_path)
        c.drawImage(
            ImageReader(io.BytesIO(jpeg)),
            0,
            0,
            width=width,
            height=height,
            preserveAspectRatio=False,
        )
        c.showPage()
    c.save()
    _strip_text_machinery(path)


def _strip_text_machinery(path: Path) -> None:
    """Remove the empty text block reportlab writes on every page.

    reportlab always emits `BT /F1 12 Tf 14.4 TL ET` and declares an unused
    Helvetica resource, even on a page that shows no text. Nothing is drawn by
    it, but a fixture that claims "no text layer" should not ship a font
    resource at all — a real scanner never produces one. Dropping both keeps
    the file honest and lets the verifier assert the strict version of the
    property.
    """
    import pikepdf

    with pikepdf.open(path, allow_overwriting_input=True) as pdf:
        for page in pdf.pages:
            raw = page.Contents.read_bytes()
            cleaned = re.sub(rb"BT\b.*?ET", b"", raw, flags=re.S)
            page.Contents = pdf.make_stream(cleaned)
            resources = page.get("/Resources")
            if resources is not None and "/Font" in resources:
                del resources["/Font"]
        # deterministic_id derives /ID from the content, not from entropy.
        pdf.save(path, deterministic_id=True)


# --------------------------------------------------------------------------
# Verification
# --------------------------------------------------------------------------


def opc_integrity(path: Path) -> None:
    """Every relationship target and content-type override must resolve.

    This is what proves ``slim_docx`` left no dangling reference behind.
    """
    with zipfile.ZipFile(path) as zf:
        names = set(zf.namelist())
        content_types = zf.read("[Content_Types].xml").decode("utf-8")
        for part in re.findall(r'PartName="/([^"]+)"', content_types):
            if part not in names:
                raise AssertionError(f"content type override for missing part: {part}")
        for rels_name in [n for n in names if n.endswith(".rels")]:
            base = os.path.dirname(os.path.dirname(rels_name))
            xml = zf.read(rels_name).decode("utf-8")
            for match in re.finditer(r'<Relationship[^>]*Target="([^"]+)"([^>]*)>', xml):
                target, rest = match.group(1), match.group(2)
                if "External" in rest or target.startswith(("http", "mailto")):
                    continue
                resolved = os.path.normpath(
                    os.path.join(base, target.lstrip("/"))
                ).replace(os.sep, "/")
                if resolved not in names:
                    raise AssertionError(
                        f"{rels_name}: relationship target missing: {target}"
                    )


def pdf_page_text(page) -> str:
    """Extract text from one page using only pikepdf.

    Walks the content stream for text-showing operators and maps the byte
    strings through the font's /ToUnicode CMap — the same map a browser PDF
    viewer uses. Returning '' therefore means "a viewer would find no text
    here", which is exactly the assertion scan.pdf needs.
    """
    import pikepdf

    cmaps = {}
    resources = page.get("/Resources")
    fonts = resources.get("/Font") if resources is not None else None
    if fonts is not None:
        for key, font in fonts.items():
            tounicode = font.get("/ToUnicode")
            if tounicode is None:
                continue
            data = tounicode.read_bytes().decode("latin-1")
            mapping = {}
            for src, dst in re.findall(r"<([0-9A-Fa-f]+)>\s*<([0-9A-Fa-f]+)>", data):
                try:
                    mapping[int(src, 16)] = bytes.fromhex(dst).decode(
                        "utf-16-be", "ignore"
                    )
                except ValueError:
                    continue
            cmaps[str(key)] = mapping

    out = []
    current = None
    for operands, operator in pikepdf.parse_content_stream(page):
        op = str(operator)
        if op == "Tf" and operands:
            current = str(operands[0])
        elif op in ("Tj", "'", '"', "TJ"):
            items = operands[-1] if op == "TJ" else [operands[-1]]
            for item in items:
                if not isinstance(item, pikepdf.String):
                    continue
                raw = bytes(item)
                mapping = cmaps.get(current) or {}
                if mapping:
                    out.append("".join(mapping.get(b, "") for b in raw))
                else:
                    out.append(raw.decode("latin-1"))
    return "".join(out)


def pdf_has_fonts(page) -> bool:
    resources = page.get("/Resources")
    if resources is None:
        return False
    fonts = resources.get("/Font")
    return fonts is not None and len(fonts) > 0


def human(size: int) -> str:
    return f"{size / 1024:.1f} KB" if size < 1024 * 1024 else f"{size / 1048576:.2f} MB"


def verify(out_dir: Path) -> int:
    """Re-open every fixture with a real parser. Returns a process exit code."""
    import docx
    import openpyxl
    import pikepdf
    from pptx import Presentation

    results = []
    failures = 0

    def record(name: str, check):
        nonlocal failures
        path = out_dir / name
        if not path.exists():
            results.append((name, "—", "MISSING"))
            failures += 1
            return
        try:
            note = check(path) or "OK"
            results.append((name, human(path.stat().st_size), note))
        except Exception as exc:  # noqa: BLE001 - report, never crash the run
            results.append(
                (name, human(path.stat().st_size), f"FAIL: {type(exc).__name__}: {exc}")
            )
            failures += 1

    def docx_chars(doc) -> int:
        chars = sum(len(p.text) for p in doc.paragraphs)
        for table in doc.tables:
            chars += sum(len(c.text) for row in table.rows for c in row.cells)
        return chars

    def check_short(path):
        doc = docx.Document(path)
        opc_integrity(path)
        assert len(doc.tables) >= 1, "expected at least one table"
        size = path.stat().st_size
        assert size < 30 * 1024, f"must stay under 30 KB, got {size}"
        pages = docx_chars(doc) // 2800 + 1
        return (
            f"docx OK · {len(doc.paragraphs)} para · {len(doc.tables)} tables · "
            f"~{pages} pages · under 30 KB"
        )

    def check_long(path):
        doc = docx.Document(path)
        opc_integrity(path)
        pages = docx_chars(doc) // 2800 + 1
        assert len(doc.tables) >= 1, "expected at least one table"
        assert pages >= 25, f"expected 25+ pages, estimated {pages}"
        size = path.stat().st_size
        assert 60 * 1024 <= size <= 200 * 1024, f"size {size} outside 60-200 KB"
        headings = sum(
            1 for p in doc.paragraphs if p.style.name.startswith(("Heading", "Title"))
        )
        return (
            f"docx OK · {len(doc.paragraphs)} para · {headings} headings · "
            f"{len(doc.tables)} tables · ~{pages} pages (est.)"
        )

    def check_heavy(path):
        doc = docx.Document(path)
        opc_integrity(path)
        size = path.stat().st_size
        assert size > 360 * 1024, f"must exceed 360 KB, got {size}"
        images = len(doc.inline_shapes)
        assert images >= 1, "expected embedded images"
        return (
            f"docx OK · {images} images · {size} B > 360 KB "
            f"(CDN cutoff {MAX_DOCX_CDN_BINARY_BYTES} B)"
        )

    def check_registry(path):
        wb = openpyxl.load_workbook(path)
        assert len(wb.sheetnames) == 3, f"expected 3 sheets, got {wb.sheetnames}"
        ws = wb[wb.sheetnames[0]]
        assert isinstance(ws["A1"].value, str) and ws["A1"].value, "A1 title missing"
        assert all(
            ws.cell(row=2, column=c).value is None for c in range(1, 9)
        ), "row 2 must be blank"
        merged = len(ws.merged_cells.ranges)
        assert merged >= 5, f"expected several merged ranges, got {merged}"
        data_rows = range(5, 45)
        empties = sum(
            1
            for r in data_rows
            for c in (5, 7, 8)
            if ws.cell(row=r, column=c).value is None
        )
        assert empties >= 10, f"expected scattered empty cells, got {empties}"
        assert all(
            isinstance(ws.cell(row=r, column=6).value, (int, float)) for r in data_rows
        ), "amounts must be real numbers, not text"
        total = ws.cell(row=45, column=6).value
        assert isinstance(total, (int, float)), "totals row must be numeric"
        ws2 = wb[wb.sheetnames[1]]
        formulas = sum(
            1
            for row in ws2.iter_rows()
            for c in row
            if isinstance(c.value, str) and c.value.startswith("=")
        )
        assert formulas >= 5, f"sheet 2 needs formulas, got {formulas}"
        # The cached formula results are injected as literal XML, so guard them
        # against drifting away from the numbers they claim to summarise.
        cached = openpyxl.load_workbook(path, data_only=True)[wb.sheetnames[1]]
        for r in range(4, 8):
            plan = ws2.cell(row=r, column=2).value
            fact = ws2.cell(row=r, column=3).value
            got = cached.cell(row=r, column=4).value
            assert got is not None, f"D{r} lost its cached value"
            assert abs(got - (fact - plan)) < 0.005, (
                f"D{r} cached {got}, but C{r}-B{r} = {fact - plan}"
            )
            assert cached.cell(row=r, column=5).value is None, (
                f"E{r} must stay uncached on purpose"
            )
        assert abs(cached["B8"].value - sum(ws2.cell(row=r, column=2).value
                                            for r in range(4, 8))) < 0.005
        assert abs(cached["C8"].value - sum(ws2.cell(row=r, column=3).value
                                            for r in range(4, 8))) < 0.005

        ws3 = wb[wb.sheetnames[2]]
        assert ws3.max_row <= 3, "sheet 3 must be nearly empty"
        return (
            f"xlsx OK · 3 sheets · {merged} merged · {empties} empty cells · "
            f"{formulas} formulas (7 cached, 4 not) · итого={total}"
        )

    def check_big_rows(path):
        wb = openpyxl.load_workbook(path, read_only=True)
        ws = wb[wb.sheetnames[0]]
        rows = ws.max_row
        wb.close()
        assert (
            rows > SPREADSHEET_MAX_ROWS_PER_SHEET
        ), f"needs >{SPREADSHEET_MAX_ROWS_PER_SHEET} rows, got {rows}"
        return f"xlsx OK · {rows} rows > {SPREADSHEET_MAX_ROWS_PER_SHEET} cutoff"

    def check_deck(expected_slides, w, h, label):
        def _check(path):
            prs = Presentation(path)
            count = len(prs.slides)
            assert count == expected_slides, (
                f"expected {expected_slides} slides, got {count}"
            )
            assert (prs.slide_width, prs.slide_height) == (w, h), (
                f"expected {w}x{h} EMU, got {prs.slide_width}x{prs.slide_height}"
            )
            titles = {
                s.shapes.title.text for s in prs.slides if s.shapes.title is not None
            }
            tables = sum(
                1 for s in prs.slides for sh in s.shapes if sh.has_table
            )
            return (
                f"pptx OK · {count} slides · {label} ({w}x{h} EMU) · "
                f"{len(titles)} distinct titles · {tables} tables"
            )

        return _check

    def check_deck_many(path):
        prs = Presentation(path)
        count = len(prs.slides)
        assert count == 60, f"expected 60 slides, got {count}"
        titles = [s.shapes.title.text for s in prs.slides if s.shapes.title is not None]
        assert len(titles) == 60, "every slide needs a title"
        assert len(set(titles)) == 60, "slide titles must be distinct"
        assert (prs.slide_width, prs.slide_height) == (12_192_000, 6_858_000)
        return "pptx OK · 60 slides · 16:9 · 60 distinct titles"

    def check_notes(path):
        text = path.read_text(encoding="utf-8")
        required = {
            "heading": r"^# ",
            "list": r"^\d+\. ",
            "table": r"^\|.*\|",
            "code fence": r"^```",
            "blockquote": r"^> ",
            "link": r"\[[^\]]+\]\(https?://[^)]+\)",
        }
        missing = [k for k, pattern in required.items() if not re.search(pattern, text, re.M)]
        assert not missing, f"markdown missing: {', '.join(missing)}"
        return f"markdown OK · {len(text.splitlines())} lines · all elements present"

    def check_script(path):
        source = path.read_text(encoding="utf-8")
        compile(source, str(path), "exec")
        lines = source.splitlines()
        assert len(lines) >= 55, f"expected ~60 lines, got {len(lines)}"
        assert '"""' in source, "expected docstrings"
        assert re.search(r"^\s*#", source, re.M), "expected comments"
        return f"python OK · compiles · {len(lines)} lines"

    def check_csv(path):
        with open(path, encoding="utf-8", newline="") as handle:
            rows = list(csv.reader(handle))
        assert len(rows) == 31, f"expected header + 30 rows, got {len(rows)}"
        flat = [cell for row in rows[1:] for cell in row]
        assert any("," in cell for cell in flat), "need a field containing a comma"
        assert any('"' in cell for cell in flat), "need a field containing a quote"
        assert any(cell == "" for cell in flat), "need an empty field"
        assert any(re.search(r"[А-Яа-я]", cell) for cell in flat), "need Cyrillic"
        return f"csv OK · {len(rows) - 1} rows · quoted/empty/Cyrillic present"

    def check_digital_pdf(path):
        with pikepdf.open(path) as pdf:
            pages = len(pdf.pages)
            assert pages == 5, f"expected 5 pages, got {pages}"
            texts = [pdf_page_text(p) for p in pdf.pages]
        for i, text in enumerate(texts, start=1):
            assert f"Раздел {i}" in text, (
                f"page {i}: heading «Раздел {i}» is not extractable"
            )
        chars = sum(len(t) for t in texts)
        return (
            f"pdf OK · 5 pages · real text layer, {chars} chars extracted · "
            "«Раздел 1»…«Раздел 5» searchable"
        )

    def check_scan_pdf(path):
        with pikepdf.open(path) as pdf:
            pages = len(pdf.pages)
            assert pages == 3, f"expected 3 pages, got {pages}"
            for i, page in enumerate(pdf.pages, start=1):
                assert not pdf_has_fonts(page), f"page {i} has font resources"
                text = pdf_page_text(page)
                assert text.strip() == "", f"page {i} leaked text: {text[:60]!r}"
                xobjects = page.get("/Resources", {}).get("/XObject")
                assert xobjects is not None and len(xobjects) >= 1, (
                    f"page {i} has no image"
                )
        return "pdf OK · 3 pages · image-only, zero extractable text"

    def check_broken(path):
        raw = path.read_bytes()
        assert raw[:4] == b"PK\x03\x04", "must keep docx magic bytes"
        assert not zipfile.is_zipfile(path), "zip must be unreadable"
        try:
            docx.Document(path)
        except Exception:
            return f"correctly unreadable · PK magic kept · {len(raw)} B truncated"
        raise AssertionError("python-docx opened the broken file")

    def check_locked_pdf(path):
        try:
            pikepdf.open(path)
        except pikepdf.PasswordError:
            pass
        else:
            raise AssertionError("opened without a password")
        with pikepdf.open(path, password=PDF_PASSWORD) as pdf:
            pages = len(pdf.pages)
        return (
            f"encrypted OK · rejects empty password · opens with «{PDF_PASSWORD}» · "
            f"{pages} pages"
        )

    def check_locked_docx(path):
        raw = path.read_bytes()
        assert raw[:8] == bytes.fromhex("d0cf11e0a1b11ae1"), (
            "expected an OLE2/CFB container"
        )
        assert "EncryptedPackage".encode("utf-16-le") in raw, (
            "no EncryptedPackage stream"
        )
        try:
            docx.Document(path)
        except Exception:
            pass
        else:
            raise AssertionError("python-docx opened the encrypted file")
        note = "encrypted OK · CFB container · python-docx rejects it"
        try:
            import msoffcrypto
        except ImportError:
            return note + " (msoffcrypto absent, structural check only)"
        with open(path, "rb") as handle:
            office = msoffcrypto.OfficeFile(handle)
            assert office.is_encrypted(), "msoffcrypto says it is not encrypted"
            office.load_key(password=DOCX_PASSWORD)
            plain = io.BytesIO()
            office.decrypt(plain)
        plain.seek(0)
        doc = docx.Document(plain)
        return note + f" · decrypts with «{DOCX_PASSWORD}» ({len(doc.paragraphs)} para)"

    def check_archive(path):
        with zipfile.ZipFile(path) as zf:
            assert zf.testzip() is None, "corrupt entries"
            names = zf.namelist()
            assert len(names) == 2, f"expected 2 entries, got {names}"
        return f"zip OK · {len(names)} text entries"

    def check_unknown(path):
        raw = path.read_bytes()
        assert not zipfile.is_zipfile(path), "must not be a zip"
        assert raw[:4] != b"%PDF", "must not be a pdf"
        return f"opaque blob OK · {len(raw)} B · no known magic"

    record("contract-short.docx", check_short)
    record("contract-long.docx", check_long)
    record("contract-heavy.docx", check_heavy)
    record("registry.xlsx", check_registry)
    record("big-rows.xlsx", check_big_rows)
    record("deck-16x9.pptx", check_deck(12, 12_192_000, 6_858_000, "16:9"))
    record("deck-4x3.pptx", check_deck(6, 9_144_000, 6_858_000, "4:3"))
    record("deck-many.pptx", check_deck_many)
    record("notes.md", check_notes)
    record("script.py", check_script)
    record("data.csv", check_csv)
    record("digital.pdf", check_digital_pdf)
    record("scan.pdf", check_scan_pdf)
    record("broken.docx", check_broken)
    record("locked.pdf", check_locked_pdf)
    if (out_dir / "locked.docx").exists():
        record("locked.docx", check_locked_docx)
    else:
        results.append(("locked.docx", "—", "SKIPPED (not generated, see README)"))
    record("archive.zip", check_archive)
    record("unknown.xyz", check_unknown)

    name_w = max(len(r[0]) for r in results) + 2
    size_w = 10
    print()
    print(f"{'FILE'.ljust(name_w)}{'SIZE'.rjust(size_w)}  RESULT")
    print("-" * (name_w + size_w + 62))
    for name, size, note in results:
        print(f"{name.ljust(name_w)}{size.rjust(size_w)}  {note}")

    files = [p for p in out_dir.iterdir() if p.is_file()]
    total = sum(p.stat().st_size for p in files)
    print("-" * (name_w + size_w + 62))
    budget = "under" if total < 2 * 1024 * 1024 else "OVER"
    print(
        f"{'TOTAL'.ljust(name_w)}{human(total).rjust(size_w)}  "
        f"{len(files)} files ({budget} the 2 MB budget)"
    )
    print(f"\n{'FAILED: ' + str(failures) if failures else 'All checks passed.'}")
    return 1 if failures else 0


# --------------------------------------------------------------------------
# Entry point
# --------------------------------------------------------------------------


def generate(out_dir: Path) -> None:
    out_dir.mkdir(parents=True, exist_ok=True)
    regular, bold = find_font()
    print(f"font (regular): {regular}")
    print(f"font (bold):    {bold}")

    print("· contract-short.docx")
    build_contract_short(out_dir / "contract-short.docx")

    print("· contract-long.docx")
    build_contract_long(
        out_dir / "contract-long.docx",
        letterhead=render_letterhead(regular),
        stamp=render_stamp(regular),
    )

    print("· contract-heavy.docx (embedding scans until it clears 350 KB)")
    scans = [
        render_scan_page(i, SCAN_PAGES[(i - 1) % len(SCAN_PAGES)], regular, quality=72)
        for i in range(1, 7)
    ]
    heavy = build_contract_heavy(out_dir / "contract-heavy.docx", scans)
    print(
        f"  -> {heavy} bytes ({heavy / 1024:.1f} KB); CDN cutoff is "
        f"{MAX_DOCX_CDN_BINARY_BYTES} bytes"
    )

    print("· registry.xlsx")
    build_registry_xlsx(out_dir / "registry.xlsx")

    print("· big-rows.xlsx")
    build_big_rows_xlsx(out_dir / "big-rows.xlsx")

    print("· deck-16x9.pptx / deck-4x3.pptx / deck-many.pptx")
    build_deck_16x9(out_dir / "deck-16x9.pptx")
    build_deck_4x3(out_dir / "deck-4x3.pptx")
    build_deck_many(out_dir / "deck-many.pptx")

    print("· notes.md / script.py / data.csv")
    (out_dir / "notes.md").write_text(NOTES_MD, encoding="utf-8")
    (out_dir / "script.py").write_text(SCRIPT_PY, encoding="utf-8")
    build_data_csv(out_dir / "data.csv")

    print("· digital.pdf / scan.pdf / locked.pdf")
    build_digital_pdf(out_dir / "digital.pdf", regular, bold)
    build_scan_pdf(out_dir / "scan.pdf", regular)
    build_locked_pdf(out_dir / "locked.pdf", regular, bold)

    print("· broken.docx / locked.docx")
    build_broken_docx(out_dir / "broken.docx")
    status = build_locked_docx(out_dir / "locked.docx")
    if status != "OK":
        print(f"  !! locked.docx {status}")
        print("  !! Nothing fake is written in its place — the fixture is absent.")

    print("· archive.zip / unknown.xyz")
    build_archive_zip(out_dir / "archive.zip")
    build_unknown_xyz(out_dir / "unknown.xyz")


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    parser.add_argument(
        "--out", type=Path, default=OUT_DIR, help="output directory (default: ./files)"
    )
    parser.add_argument(
        "--verify",
        action="store_true",
        help="only re-check existing fixtures, write nothing",
    )
    args = parser.parse_args()

    if not args.verify:
        generate(args.out)
    return verify(args.out)


if __name__ == "__main__":
    sys.exit(main())
