#!/usr/bin/env python3
"""Check that the committed fixtures are what the e2e specs assume they are.

Run after `generate_fixtures.py`. Every claim the specs rely on is asserted
here — a valid file must open, a damaged one must fail to open, the scan must
carry no extractable text, the oversized workbook must really exceed the row
cap. Without this the fixtures could drift into something that still uploads
but no longer exercises the path the test is named after.
"""

from __future__ import annotations

import sys
import zipfile
from pathlib import Path

import pikepdf
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

FILES = Path(__file__).resolve().parent / "files"

CDN_BOUND_BYTES = 350 * 1024
ROW_CAP = 5000

problems: list[str] = []
report: list[tuple[str, int, str]] = []


def check(name: str, note: str, condition: bool, detail: str = "") -> None:
    size = (FILES / name).stat().st_size if (FILES / name).exists() else 0
    report.append((name, size, note if condition else f"FAILED: {detail or note}"))
    if not condition:
        problems.append(f"{name}: {detail or note}")


def docx_text(path: Path) -> str:
    document = Document(path)
    parts = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            parts.extend(cell.text for cell in row.cells)
    return "\n".join(parts)


def pdf_text(path: Path, password: str = "") -> str:
    with pikepdf.open(path, password=password) as pdf:
        extracted = []
        for page in pdf.pages:
            content = page.obj.get("/Contents")
            if content is None:
                continue
            data = bytes(content.read_bytes()) if hasattr(content, "read_bytes") else b""
            extracted.append(data.decode("latin-1", errors="ignore"))
    return "\n".join(extracted)


def main() -> int:
    short_text = docx_text(FILES / "contract-short.docx")
    check(
        "contract-short.docx",
        "opens, carries contract text",
        "Договор" in short_text and "Ромашка" in short_text,
    )

    long_document = Document(FILES / "contract-long.docx")
    check(
        "contract-long.docx",
        f"opens, {len(long_document.paragraphs)} paragraphs",
        len(long_document.paragraphs) > 100,
        "expected a genuinely long document",
    )

    heavy = (FILES / "contract-heavy.docx").stat().st_size
    check(
        "contract-heavy.docx",
        f"opens, {heavy:,} bytes — above the {CDN_BOUND_BYTES:,} bound",
        heavy > CDN_BOUND_BYTES + 10_000 and bool(docx_text(FILES / "contract-heavy.docx")),
        f"must exceed {CDN_BOUND_BYTES:,} bytes with margin, got {heavy:,}",
    )

    registry = load_workbook(FILES / "registry.xlsx")
    merged = len(registry["Реестр договоров"].merged_cells.ranges)
    check(
        "registry.xlsx",
        f"opens, sheets={registry.sheetnames}, merged ranges={merged}",
        len(registry.sheetnames) == 3 and merged >= 2,
        "expected three sheets and merged cells",
    )

    big = load_workbook(FILES / "big-rows.xlsx", read_only=True)
    rows = big["Операции"].max_row
    check(
        "big-rows.xlsx",
        f"opens, {rows:,} rows — above the {ROW_CAP:,} cap",
        rows > ROW_CAP,
        f"expected more than {ROW_CAP} rows, got {rows}",
    )

    for name, expected_ratio, expected_slides in (
        ("deck-16x9.pptx", 16 / 9, 12),
        ("deck-4x3.pptx", 4 / 3, 6),
        ("deck-many.pptx", 16 / 9, 60),
    ):
        deck = Presentation(FILES / name)
        ratio = deck.slide_width / deck.slide_height
        check(
            name,
            f"opens, {len(deck.slides._sldIdLst)} slides, ratio {ratio:.3f}",
            abs(ratio - expected_ratio) < 0.01 and len(deck.slides._sldIdLst) == expected_slides,
            f"expected {expected_slides} slides at ratio {expected_ratio:.3f}, got "
            f"{len(deck.slides._sldIdLst)} at {ratio:.3f}",
        )

    digital = pdf_text(FILES / "digital.pdf")
    with pikepdf.open(FILES / "digital.pdf") as pdf:
        digital_pages = len(pdf.pages)
    check(
        "digital.pdf",
        f"opens, {digital_pages} pages with a text layer",
        digital_pages == 5 and "TJ" in digital or "Tj" in digital,
        "expected 5 pages carrying text-drawing operators",
    )

    scan = pdf_text(FILES / "scan.pdf")
    with pikepdf.open(FILES / "scan.pdf") as pdf:
        scan_pages = len(pdf.pages)
    check(
        "scan.pdf",
        f"opens, {scan_pages} image-only pages",
        scan_pages == 3 and "Tj" not in scan and "TJ" not in scan,
        "scan must have no text-drawing operators at all",
    )

    opened_without_password = True
    try:
        with pikepdf.open(FILES / "locked.pdf"):
            pass
    except pikepdf.PasswordError:
        opened_without_password = False
    with pikepdf.open(FILES / "locked.pdf", password="secret123") as pdf:
        unlocked_pages = len(pdf.pages)
    check(
        "locked.pdf",
        f"refuses without a password, opens with one ({unlocked_pages} pages)",
        not opened_without_password and unlocked_pages == 5,
        "expected the file to require its password",
    )

    broken_opens = True
    try:
        Document(FILES / "broken.docx")
    except Exception:
        broken_opens = False
    check(
        "broken.docx",
        "is genuinely unreadable",
        not broken_opens,
        "a damaged fixture that still parses proves nothing",
    )

    with zipfile.ZipFile(FILES / "archive.zip") as archive:
        names = archive.namelist()
    check("archive.zip", f"valid archive containing {names}", len(names) == 2)

    for name, marker in (
        ("notes.md", "## Что сделано"),
        ("script.py", "def total_amount"),
        ("data.csv", "Контрагент"),
    ):
        text = (FILES / name).read_text(encoding="utf-8")
        check(name, "readable text fixture", marker in text, f"missing marker {marker!r}")

    total = sum(item.stat().st_size for item in FILES.iterdir())
    width = max(len(name) for name, _, _ in report)
    for name, size, note in report:
        print(f"{name:<{width}}  {size:>10,}  {note}")
    print(f"\ntotal {total:,} bytes across {len(report)} fixtures")

    if problems:
        print(f"\n{len(problems)} problem(s):")
        for problem in problems:
            print(f"  {problem}")
        return 1
    return 0


if __name__ == "__main__":
    sys.exit(main())
