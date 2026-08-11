#!/usr/bin/env python3
"""Run the Russian-first PPTX golden matrix in the real office toolchain."""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import shutil
import subprocess
import sys
import tempfile
import zipfile
from copy import deepcopy
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[2]
BENCH_DIR = Path(__file__).resolve().parent
BUILDER = ROOT / "skill/pptx/scripts/build_presentation.py"
BASE_FIXTURE = BENCH_DIR / "fixtures/ru_board.json"
CASES_PATH = BENCH_DIR / "cases.json"
SCORES_PATH = BENCH_DIR / "visual_scores.json"
RESULTS_DIR = BENCH_DIR / "results"

SOURCE = "Источник: управленческий отчёт за июль 2026"
SOURCES = [
    {"label": "Управленческий отчёт", "url": "management-report.xlsx — июль 2026"},
    {"label": "CRM-выгрузка", "url": "crm-export.csv — срез от 31.07.2026"},
]
RENDER_EVIDENCE_ALGORITHM = "sha256-json-v1"


def build_render_evidence_digest(
    results: list[dict[str, Any]], expected_cases: list[str]
) -> tuple[str | None, list[str]]:
    """Bind a human visual scorecard to the exact freshly rendered matrix.

    The full-slide hashes are intentionally platform-specific. A renderer or
    layout change must therefore invalidate the recorded human review instead
    of silently reusing an old 9/10 score.
    """
    evidence: dict[str, list[str]] = {}
    for result in results:
        if result.get("run") != 1 or result.get("case") not in expected_cases:
            continue
        image_hashes = result.get("imageHashes")
        if isinstance(image_hashes, list) and image_hashes and all(
            isinstance(value, str) and value for value in image_hashes
        ):
            evidence[str(result["case"])] = image_hashes

    missing = [case_id for case_id in expected_cases if case_id not in evidence]
    if missing:
        return None, [
            "render evidence is missing for: " + ", ".join(sorted(missing))
        ]

    canonical = json.dumps(
        {case_id: evidence[case_id] for case_id in sorted(expected_cases)},
        ensure_ascii=False,
        separators=(",", ":"),
    )
    return hashlib.sha256(canonical.encode("utf-8")).hexdigest(), []


def evaluate_visual_scores(scores: dict[str, float]) -> list[str]:
    """Return acceptance issues for the 9/10 professional visual gate."""
    if not scores:
        return ["visual scores are missing"]
    values = list(scores.values())
    issues: list[str] = []
    average = sum(values) / len(values)
    if average < 9.0:
        issues.append(f"visual average is {average:.2f}, below 9.0")
    if any(score < 8.0 for score in values):
        issues.append("one or more visual scores are below 8")
    professional = sum(score >= 9.0 for score in values) / len(values)
    if professional < 0.9:
        issues.append(f"only {professional:.0%} of visual scores are at least 9.0")
    return issues


def evaluate_visual_scorecard(
    payload: dict[str, Any], expected_cases: list[str], render_evidence: str | None
) -> list[str]:
    """Validate reviewer coverage, rubric arithmetic, and the visual quality gate."""
    rubric = payload.get("rubric", {})
    cases = payload.get("cases", {})
    issues: list[str] = []
    recorded_evidence = payload.get("renderEvidence")
    if render_evidence is None:
        issues.append("current render evidence is unavailable")
    elif not isinstance(recorded_evidence, dict):
        issues.append("visual scorecard render evidence is missing")
    elif recorded_evidence.get("algorithm") != RENDER_EVIDENCE_ALGORITHM:
        issues.append("visual scorecard render evidence uses an unsupported algorithm")
    elif recorded_evidence.get("matrixDigest") != render_evidence:
        issues.append(
            "visual scorecard evidence is stale; review the fresh renders and update the scorecard"
        )
    scores: dict[str, float] = {}
    for case_id in expected_cases:
        value = cases.get(case_id, {})
        if value.get("score") is None:
            issues.append(f"visual score is missing for {case_id}")
            continue
        score = float(value["score"])
        dimensions = value.get("dimensions", {})
        if set(dimensions) != set(rubric):
            issues.append(f"visual dimensions are incomplete for {case_id}")
        else:
            for dimension, dimension_score in dimensions.items():
                if not 0 <= float(dimension_score) <= float(rubric[dimension]):
                    issues.append(
                        f"visual dimension {dimension} is out of range for {case_id}"
                    )
            if abs(sum(float(item) for item in dimensions.values()) - score) > 0.01:
                issues.append(f"visual score does not match dimension sum for {case_id}")
        scores[case_id] = score
    if set(cases) != set(expected_cases):
        issues.append("visual scorecard case coverage does not match the golden matrix")
    issues.extend(evaluate_visual_scores(scores))
    return issues


def _job(case_id: str, filename: str) -> dict[str, Any]:
    return {
        "format": "pptx",
        "audience": "Руководители компании",
        "goal": f"Принять решение по сценарию {case_id}",
        "sourceFileIds": ["management-report", "crm-export"],
        "immutableElements": ["Фактические значения"],
        "locale": "ru-RU",
        "filename": filename,
        "acceptanceCriteria": [
            "Кириллица отображается корректно",
            "Фактологические слайды имеют источники",
            "Диаграммы и схемы остаются редактируемыми",
        ],
    }


def _title(title: str, subtitle: str) -> dict[str, Any]:
    return {
        "layout": "title",
        "eyebrow": "МАТЕРИАЛЫ К РЕШЕНИЮ",
        "title": title,
        "subtitle": subtitle,
        "date": "Август 2026",
    }


def _claim(title: str, claim: str, support: str) -> dict[str, Any]:
    return {
        "layout": "claim",
        "title": title,
        "claim": claim,
        "support": support,
        "source": SOURCE,
    }


def _summary(*items: str) -> dict[str, Any]:
    return {"layout": "summary", "title": "Следующие решения", "bullets": list(items)}


def _base_spec(case_id: str, filename: str) -> dict[str, Any]:
    return {
        "job": _job(case_id, filename),
        "slides": [],
        "sources": deepcopy(SOURCES),
        "changeLog": [
            {"target": "Presentation", "summary": f"Created golden scenario {case_id}"}
        ],
        "repairIterations": 0,
        "outputPdf": True,
    }


def _case_spec(case_id: str, filename: str, asset_dir: Path) -> dict[str, Any]:
    if case_id == "ru-board-notes":
        spec = json.loads(BASE_FIXTURE.read_text(encoding="utf-8"))
        spec["job"]["filename"] = filename
        return spec

    spec = _base_spec(case_id, filename)
    if case_id == "ru-sales-proposal":
        image_path = asset_dir / "implementation-map.png"
        _make_illustration(image_path)
        spec["slides"] = [
            _title("Управляемое внедрение за шесть недель", "Предложение для директора по закупкам"),
            _claim(
                "Ценность появляется до завершения полной интеграции",
                "Первая команда начинает работать уже на четвёртой неделе",
                "План разделяет обязательный запуск и последующее масштабирование.",
            ),
            {
                "layout": "image",
                "title": "Контур внедрения охватывает людей, данные и контроль",
                "image": str(image_path),
                "caption": "Схема целевого контура внедрения",
                "source": "Источник: проектный план внедрения v3",
            },
            {
                "layout": "comparison",
                "title": "Предложение снижает риск без потери скорости",
                "left": {"heading": "Типовой проект", "bullets": ["Единый большой запуск", "Риск выявляется поздно", "Оплата по календарю"]},
                "right": {"heading": "Предлагаемый проект", "bullets": ["Две независимые волны", "Контроль каждую неделю", "Оплата по результатам"]},
                "source": "Источник: проектный план и условия предложения",
            },
            {
                "layout": "process",
                "title": "Шесть недель разбиты на четыре контрольные точки",
                "steps": [
                    {"title": "Диагностика", "detail": "Неделя 1"},
                    {"title": "Настройка", "detail": "Недели 2–3"},
                    {"title": "Пилот", "detail": "Неделя 4"},
                    {"title": "Масштаб", "detail": "Недели 5–6"},
                ],
                "source": "Источник: проектный план внедрения v3",
            },
            _summary("Согласовать объём пилота", "Назначить владельцев данных", "Утвердить дату старта"),
        ]
    elif case_id == "ru-xlsx-analysis":
        spec["slides"] = [
            _title("Финансовая модель роста", "Тезисы для инвестиционного комитета"),
            _claim("Базовый сценарий финансируется из операционного потока", "Рост не требует внешнего капитала", "Запас ликвидности остаётся выше внутреннего минимума."),
            _chart("Выручка растёт быстрее после запуска канала", [82, 114, 161], "+41%", "2027 к 2026"),
            _chart("Валовая прибыль догоняет рост выручки", [31, 44, 67], "+23 млн ₽", "дополнительная прибыль"),
            {
                "layout": "table",
                "title": "Базовый сценарий сохраняет запас ликвидности",
                "columns": ["Показатель", "2025", "2026", "2027"],
                "rows": [["Денежный поток, млн ₽", "18", "24", "39"], ["Запас, мес.", "6,1", "5,8", "7,2"], ["CAPEX, млн ₽", "7", "11", "14"]],
                "source": "Источник: финансовая модель v7, лист «ДДС»",
            },
            _metrics(),
            _summary("Утвердить базовый сценарий", "Зафиксировать лимит CAPEX", "Обновлять модель ежемесячно"),
        ]
    elif case_id == "ru-docx-strategy":
        spec["slides"] = [
            _title("Стратегия корпоративного направления", "Решения на горизонте двенадцати месяцев"),
            {"layout": "section", "number": "01", "title": "Где играть", "subtitle": "Сегменты с коротким путём к доказанной ценности"},
            _claim("Сосредоточение на двух отраслях повышает конверсию", "Ритейл и логистика дают 68% подтверждённой воронки", "В этих сегментах уже есть повторяемые сценарии внедрения."),
            {
                "layout": "comparison",
                "title": "Две отрасли требуют разных аргументов, но общей платформы",
                "left": {"heading": "Ритейл", "bullets": ["Скорость операций", "Распределённые команды", "Эффект за 30 дней"]},
                "right": {"heading": "Логистика", "bullets": ["Контроль отклонений", "Интеграция данных", "Эффект за 45 дней"]},
                "source": "Источник: стратегическое мемо, раздел 3",
            },
            {
                "layout": "process",
                "title": "Стратегия переводится в квартальный цикл исполнения",
                "steps": [{"title": "Выбор", "detail": "Q3"}, {"title": "Доказательство", "detail": "Q4"}, {"title": "Масштаб", "detail": "Q1"}, {"title": "Оптимизация", "detail": "Q2"}],
                "source": "Источник: стратегическое мемо, дорожная карта",
            },
            _summary("Сфокусировать маркетинг на двух отраслях", "Собрать отраслевые пакеты", "Проверять стратегию ежеквартально"),
        ]
    elif case_id == "ru-long-copy":
        spec["slides"] = [
            _title("Как вырасти в корпоративном сегменте без потери качества внедрения", "План решений для руководителей функций"),
            _claim("Ключевое ограничение находится не в продажах, а в скорости получения первой измеримой пользы", "Рост зависит от первых четырёх недель работы клиента", "Поэтому ресурсы нужно перенести из поздней поддержки в раннее внедрение."),
            {
                "layout": "two_column",
                "title": "Перераспределение ответственности сокращает путь от подписания договора до первой ценности",
                "left": {"heading": "До изменения", "bullets": ["Менеджер координирует все команды", "Интеграции начинаются после обучения", "Риски обсуждаются на финальном этапе"]},
                "right": {"heading": "После изменения", "bullets": ["Владелец внедрения назначается заранее", "Интеграции и обучение идут параллельно", "Риски фиксируются в первую неделю"]},
                "source": "Источник: анализ 32 корпоративных запусков",
            },
            _chart("Сокращение срока внедрения повышает долю клиентов, дошедших до регулярного использования", [52, 67, 81], "+29 п.п.", "между крайними группами"),
            _metrics(),
            _summary("Назначить владельца внедрения", "Запустить параллельный поток интеграций", "Контролировать первую ценность"),
        ]
    elif case_id == "ru-source-traceability":
        spec["slides"] = [
            _title("Доказательная база плана роста", "Факты, допущения и решения"),
            _claim("Подтверждённая воронка покрывает базовый сценарий", "74% плана уже связано с конкретными возможностями", "Оставшаяся часть обеспечивается исторической конверсией текущего сегмента."),
            _metrics(),
            _chart("Прогноз опирается на наблюдаемую динамику воронки", [61, 79, 103], "+30%", "год к году"),
            {
                "layout": "table",
                "title": "Каждое допущение имеет владельца и контрольную дату",
                "columns": ["Допущение", "Владелец", "Проверка"],
                "rows": [["Конверсия 28%", "Продажи", "15.09.2026"], ["Срок 28 дней", "Внедрение", "30.09.2026"], ["Отток 4,2%", "CS", "Ежемесячно"]],
                "source": "Источник: реестр допущений v4",
            },
            _summary("Утвердить допущения", "Назначить контрольные даты", "Показать отклонения на следующем обзоре"),
        ]
    elif case_id == "ru-process-deck":
        spec["slides"] = [
            _title("Новый процесс запуска продукта", "Единый ритм от идеи до масштаба"),
            _claim("Решение о масштабе принимается после измеримого пилота", "Команда отделяет проверку спроса от промышленного запуска", "Это уменьшает стоимость ошибки и ускоряет сильные инициативы."),
            {
                "layout": "process",
                "title": "Инициатива проходит пять прозрачных ворот",
                "steps": [{"title": "Проблема", "detail": "1 неделя"}, {"title": "Гипотеза", "detail": "1 неделя"}, {"title": "Пилот", "detail": "3 недели"}, {"title": "Решение", "detail": "1 день"}, {"title": "Масштаб", "detail": "6 недель"}],
                "source": "Источник: продуктовый регламент v2",
            },
            {
                "layout": "process",
                "title": "Каждые ворота требуют одного проверяемого результата",
                "steps": [{"title": "Интервью", "detail": "Проблема"}, {"title": "Прототип", "detail": "Гипотеза"}, {"title": "Метрика", "detail": "Пилот"}, {"title": "Решение", "detail": "Комитет"}],
                "source": "Источник: продуктовый регламент v2",
            },
            {
                "layout": "comparison",
                "title": "Новый ритм убирает незаметные долгие проекты",
                "left": {"heading": "Раньше", "bullets": ["Срок без ограничения", "Решение по ощущениям", "Нет владельца эффекта"]},
                "right": {"heading": "Теперь", "bullets": ["Лимит шесть недель", "Единая метрика пилота", "Владелец назначен заранее"]},
                "source": "Источник: продуктовый регламент v2",
            },
            _summary("Утвердить новые ворота", "Назначить состав комитета", "Запустить процесс с сентября"),
        ]
    elif case_id == "ru-en-mixed-business":
        spec["slides"] = [
            _title("План роста ARR на 2027 год", "Бизнес-метрики, бюджет и контрольные даты"),
            _claim("Рост ARR финансируется без ухудшения unit-экономики", "LTV/CAC остаётся выше 4,0×", "Бюджет увеличивается только после подтверждения payback на пилотной когорте."),
            {
                "layout": "metrics",
                "title": "Ключевые метрики сохраняют запас прочности",
                "metrics": [{"value": "₽161 млн", "label": "ARR", "detail": "+41% YoY"}, {"value": "4,3×", "label": "LTV/CAC", "detail": "цель ≥ 4,0×"}, {"value": "7 мес.", "label": "Payback", "detail": "на 31.07.2026"}],
                "source": "Источник: SaaS model v7, вкладка «Unit economics»",
            },
            _chart("ARR ускоряется при сохранении целевого Payback", [82, 114, 161], "+₽47 млн", "прирост за год"),
            {
                "layout": "table",
                "title": "Контрольные точки привязаны к датам и метрикам",
                "columns": ["Milestone", "Дата", "Gate"],
                "rows": [["Pilot", "15.09.2026", "Payback ≤ 8 мес."], ["Scale", "01.11.2026", "NRR ≥ 108%"], ["Review", "15.01.2027", "ARR ≥ ₽125 млн"]],
                "source": "Источник: SaaS model v7 и operating plan",
            },
            _summary("Утвердить бюджет ₽18 млн", "Проверить Pilot 15.09.2026", "Масштабировать только после Gate"),
        ]
    else:
        raise ValueError(f"No net-new golden spec for {case_id}")
    return spec


def _chart(title: str, values: list[int], takeaway: str, detail: str) -> dict[str, Any]:
    return {
        "layout": "chart",
        "title": title,
        "chart": {"type": "column", "categories": ["2025", "2026", "2027"], "series": [{"name": "Значение", "values": values}], "numberFormat": "0"},
        "takeaway": takeaway,
        "detail": detail,
        "source": "Источник: финансовая модель v7",
    }


def _metrics() -> dict[str, Any]:
    return {
        "layout": "metrics",
        "title": "Экономика подтверждает возможность масштабирования",
        "metrics": [{"value": "41%", "label": "рост", "detail": "2027 к 2026"}, {"value": "4,2%", "label": "отток", "detail": "последние 12 месяцев"}, {"value": "7 мес.", "label": "окупаемость", "detail": "базовый сценарий"}],
        "source": SOURCE,
    }


def _make_illustration(path: Path) -> None:
    image = Image.new("RGB", (1600, 900), "#181F2A")
    draw = ImageDraw.Draw(image)
    colors = ["#E25843", "#F8F9FA", "#25805C"]
    for index, x in enumerate((310, 800, 1290)):
        draw.ellipse((x - 105, 345, x + 105, 555), fill=colors[index])
    draw.line((415, 450, 695, 450), fill="#DCE0E6", width=18)
    draw.line((905, 450, 1185, 450), fill="#DCE0E6", width=18)

    # People, data, and control pictograms keep the image meaningful without
    # introducing font dependencies into the deterministic golden fixture.
    for x in (278, 310, 342):
        draw.ellipse((x - 13, 399, x + 13, 425), fill="#FFFFFF")
    draw.rounded_rectangle((257, 436, 363, 495), radius=20, fill="#FFFFFF")
    for y in (397, 434, 471):
        draw.ellipse((744, y, 856, y + 32), outline="#181F2A", width=8)
        if y < 471:
            draw.rectangle((744, y + 16, 856, y + 37), fill="#F8F9FA")
            draw.line((744, y + 37, 856, y + 37), fill="#181F2A", width=8)
    draw.line((1242, 451, 1278, 487), fill="#FFFFFF", width=16)
    draw.line((1277, 487, 1342, 408), fill="#FFFFFF", width=16)
    image.save(path)


def _template_spec(case_id: str, filename: str, template_path: Path) -> dict[str, Any]:
    template = Presentation()
    template.save(template_path)
    layouts = {layout.name.lower(): layout.name for layout in template.slide_layouts}
    title_layout = layouts.get("title slide", template.slide_layouts[0].name)
    body_layout = layouts.get("title and content", template.slide_layouts[1].name)
    two_layout = layouts.get("two content", body_layout)
    spec = _base_spec(case_id, filename)
    spec["templatePath"] = str(template_path)
    spec["slides"] = [
        {**_title("План роста в корпоративном шаблоне", "Проверка мастера, макетов и редактируемости"), "templateLayout": title_layout},
        {**_claim("Шаблон остаётся основой каждой страницы", "Контент использует унаследованные поля", "Мастер, тема и геометрия не заменяются новым оформлением."), "templateLayout": body_layout},
        {"layout": "bullets", "templateLayout": body_layout, "title": "Три принципа решения", "bullets": ["Сохранить мастер", "Заполнить поля", "Не менять геометрию"]},
        {"layout": "comparison", "templateLayout": two_layout, "title": "Содержание меняется, система шаблона остаётся", "left": {"heading": "Сохраняем", "bullets": ["Мастер", "Макеты", "Размер"]}, "right": {"heading": "Заполняем", "bullets": ["Заголовки", "Тезисы", "Источники"]}, "source": "Источник: требования к корпоративному шаблону"},
        {**_summary("Утвердить структуру", "Проверить бренд-командой", "Сохранить как новую версию"), "templateLayout": body_layout},
        {"layout": "sources", "templateLayout": body_layout, "title": "Источники", "entries": deepcopy(SOURCES)},
    ]
    return spec


def _run_builder(spec: dict[str, Any], output: Path) -> dict[str, Any]:
    spec_path = output.with_suffix(".spec.json")
    spec_path.write_text(json.dumps(spec, ensure_ascii=False, indent=2), encoding="utf-8")
    result = subprocess.run(
        [sys.executable, str(BUILDER), str(spec_path), str(output)],
        cwd=ROOT,
        capture_output=True,
        text=True,
        timeout=120,
        check=False,
        env=os.environ.copy(),
    )
    if result.returncode != 0:
        raise RuntimeError((result.stderr or result.stdout).strip())
    return json.loads(Path(f"{output}.artifact-report.json").read_text(encoding="utf-8"))


def _reset_run_dir(run_dir: Path) -> None:
    if run_dir.exists():
        shutil.rmtree(run_dir)
    run_dir.mkdir(parents=True)


def _write_montage(slides: list[Path], output: Path) -> None:
    columns = min(5, max(len(slides), 1))
    tile_width, tile_height, gutter = 480, 270, 18
    rows = (len(slides) + columns - 1) // columns
    montage = Image.new(
        "RGB",
        (
            columns * tile_width + (columns + 1) * gutter,
            rows * tile_height + (rows + 1) * gutter,
        ),
        "#E6E6E6",
    )
    for index, slide_path in enumerate(slides):
        with Image.open(slide_path) as slide_image:
            tile = slide_image.convert("RGB")
            tile.thumbnail((tile_width, tile_height))
            column = index % columns
            row = index // columns
            x = gutter + column * (tile_width + gutter) + (tile_width - tile.width) // 2
            y = gutter + row * (tile_height + gutter) + (tile_height - tile.height) // 2
            montage.paste(tile, (x, y))
    montage.save(output)


def _render_review_assets(output: Path) -> dict[str, Any]:
    """Persist current full-slide renders so a reviewer never sees stale images."""
    pdf = output.with_suffix(".pdf")
    rasterizer = shutil.which("pdftoppm")
    if not pdf.exists() or not rasterizer:
        return {}
    slides_dir = output.parent / "slides"
    if slides_dir.exists():
        shutil.rmtree(slides_dir)
    slides_dir.mkdir(parents=True)
    prefix = slides_dir / "slide"
    result = subprocess.run(
        [rasterizer, "-png", "-r", "120", str(pdf), str(prefix)],
        capture_output=True,
        text=True,
        timeout=60,
        check=False,
    )
    if result.returncode != 0:
        raise RuntimeError((result.stderr or result.stdout).strip())
    slides = sorted(
        slides_dir.glob("slide-*.png"),
        key=lambda path: int(path.stem.rsplit("-", 1)[-1]),
    )
    if not slides:
        raise RuntimeError("PDF rasterizer produced no review images")
    montage = output.parent / "montage.png"
    _write_montage(slides, montage)
    return {
        "slides": [str(path.relative_to(ROOT)) for path in slides],
        "montage": str(montage.relative_to(ROOT)),
    }


def _slide_xml(path: Path, number: int) -> bytes:
    with zipfile.ZipFile(path) as archive:
        return archive.read(f"ppt/slides/slide{number}.xml")


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _validate_requirements(
    case: dict[str, Any], output: Path, report: dict[str, Any], context: dict[str, Any]
) -> list[str]:
    requirements = case["requirements"]
    deck = Presentation(output)
    issues: list[str] = []
    if report.get("status") != "ready":
        issues.append("artifact report is not ready")
    if any(issue.get("severity") == "critical" for issue in report.get("issues", [])):
        issues.append("artifact report contains a critical issue")
    if len(deck.slides) < int(requirements.get("minSlides", 0)):
        issues.append(f"deck has only {len(deck.slides)} slides")
    native_charts = sum(
        shape.shape_type == MSO_SHAPE_TYPE.CHART
        for slide in deck.slides
        for shape in slide.shapes
    )
    if native_charts < int(requirements.get("nativeCharts", 0)):
        issues.append(f"deck has only {native_charts} native charts")
    all_text = " ".join(
        shape.text
        for slide in deck.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    if requirements.get("sources") and "Источ" not in all_text:
        issues.append("sources slide is missing")
    if requirements.get("editableProcess") and "Process step title" not in {
        shape.name for slide in deck.slides for shape in slide.shapes
    }:
        issues.append("editable process shapes are missing")
    if requirements.get("mixedTypography") and not all(token in all_text for token in ("₽", "%", "ARR")):
        issues.append("mixed Russian business typography is incomplete")
    if requirements.get("longRussianCopy") and not any(
        len(shape.text) >= 70
        for slide in deck.slides
        for shape in slide.shapes
        if shape.name == "Slide title" and getattr(shape, "has_text_frame", False)
    ):
        issues.append("long Russian title was not exercised")
    if requirements.get("immutableInput") and context.get("sourceHashBefore") != context.get("sourceHashAfter"):
        issues.append("source PPTX was modified")
    if requirements.get("targetedEdit"):
        source = context["source"]
        for number in context["unaffectedSlides"]:
            if _slide_xml(source, number) != _slide_xml(output, number):
                issues.append(f"unaffected slide {number} changed")
    if requirements.get("templateFidelity"):
        source = Presentation(context["template"])
        same = (
            len(source.slide_masters) == len(deck.slide_masters)
            and sum(len(master.slide_layouts) for master in source.slide_masters)
            == sum(len(master.slide_layouts) for master in deck.slide_masters)
            and source.slide_width == deck.slide_width
            and source.slide_height == deck.slide_height
        )
        if not same:
            issues.append("template masters, layouts, or geometry changed")
    return issues


def _prepare_case(case_id: str, run_dir: Path, filename: str) -> tuple[dict[str, Any], dict[str, Any]]:
    if case_id == "ru-template-fidelity":
        template = run_dir / "template.pptx"
        return _template_spec(case_id, filename, template), {"template": template}
    if case_id == "ru-targeted-slide-edit":
        source = run_dir / "source.pptx"
        source_spec = json.loads(BASE_FIXTURE.read_text(encoding="utf-8"))
        source_spec["job"]["filename"] = source.name
        source_spec["outputPdf"] = False
        _run_builder(source_spec, source)
        before = _sha256(source)
        spec = {
            "job": _job(case_id, filename),
            "inputPath": str(source),
            "edits": [{"slide": 4, "replacements": {"+47 млн ₽": "+52 млн ₽"}, "summary": "Updated the approved revenue forecast"}],
            "outputPdf": True,
        }
        return spec, {"source": source, "sourceHashBefore": before, "unaffectedSlides": [1, 2, 3, 5, 6, 7, 8]}
    return _case_spec(case_id, filename, run_dir), {}


def run_matrix(runs: int, selected_case: str | None, skip_visual_gate: bool) -> int:
    cases = json.loads(CASES_PATH.read_text(encoding="utf-8"))
    if selected_case:
        cases = [case for case in cases if case["id"] == selected_case]
        if not cases:
            raise ValueError(f"Unknown case: {selected_case}")
    results: list[dict[str, Any]] = []
    for case in cases:
        hashes: list[list[str]] = []
        for run_number in range(1, runs + 1):
            run_dir = RESULTS_DIR / case["id"] / f"run-{run_number}"
            _reset_run_dir(run_dir)
            output = run_dir / f"{case['id']}.pptx"
            spec, context = _prepare_case(case["id"], run_dir, output.name)
            report = _run_builder(spec, output)
            review_assets = _render_review_assets(output)
            if "source" in context:
                context["sourceHashAfter"] = _sha256(context["source"])
            issues = _validate_requirements(case, output, report, context)
            raster = next(
                (check for check in report["qaChecks"] if check["name"] == "visual-raster"),
                {},
            )
            image_hashes = raster.get("details", {}).get("imageHashes", [])
            hashes.append(image_hashes)
            results.append(
                {
                    "case": case["id"],
                    "run": run_number,
                    "status": "passed" if not issues else "failed",
                    "issues": issues,
                    "artifact": str(output.relative_to(ROOT)),
                    "imageHashes": image_hashes,
                    "reviewAssets": review_assets,
                }
            )
        if not hashes[0] or any(value != hashes[0] for value in hashes[1:]):
            results.append(
                {
                    "case": case["id"],
                    "run": "determinism",
                    "status": "failed",
                    "issues": ["rendered slide pixels differ between repeated runs"],
                }
            )

    expected_case_ids = [str(case["id"]) for case in cases]
    render_evidence, render_evidence_issues = build_render_evidence_digest(
        results, expected_case_ids
    )
    visual_issues: list[str] = [*render_evidence_issues]
    if not skip_visual_gate and not selected_case:
        score_payload = json.loads(SCORES_PATH.read_text(encoding="utf-8"))
        visual_issues.extend(
            evaluate_visual_scorecard(score_payload, expected_case_ids, render_evidence)
        )

    failed = [result for result in results if result["status"] == "failed"]
    summary = {
        "runsPerCase": runs,
        "renderEvidence": {
            "algorithm": RENDER_EVIDENCE_ALGORITHM,
            "matrixDigest": render_evidence,
        },
        "results": results,
        "visualGateIssues": visual_issues,
        "status": "passed" if not failed and not visual_issues else "failed",
    }
    RESULTS_DIR.mkdir(parents=True, exist_ok=True)
    (RESULTS_DIR / "matrix-summary.json").write_text(
        json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    print(json.dumps(summary, ensure_ascii=False, indent=2))
    return 0 if summary["status"] == "passed" else 1


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--runs", type=int, default=3)
    parser.add_argument("--case")
    parser.add_argument("--skip-visual-gate", action="store_true")
    args = parser.parse_args()
    if args.runs < 1:
        parser.error("--runs must be positive")
    return run_matrix(args.runs, args.case, args.skip_visual_gate)


if __name__ == "__main__":
    raise SystemExit(main())
