#!/usr/bin/env python3
"""Structural acceptance tests for the Russian-first presentation builder."""

from __future__ import annotations

import importlib.util
import json
import sys
import tempfile
import unittest
import zipfile
from collections import Counter
from pathlib import Path
from unittest import mock

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[2]
BUILDER_PATH = ROOT / "skill/pptx/scripts/build_presentation.py"


def _load_builder():
    module_spec = importlib.util.spec_from_file_location("pptx_builder", BUILDER_PATH)
    if module_spec is None or module_spec.loader is None:
        raise RuntimeError(f"Cannot import {BUILDER_PATH}")
    module = importlib.util.module_from_spec(module_spec)
    sys.modules[module_spec.name] = module
    module_spec.loader.exec_module(module)
    return module


BUILDER = _load_builder()


def _job(filename: str = "presentation.pptx") -> dict:
    return {
        "format": "pptx",
        "audience": "Руководители компании",
        "goal": "Согласовать план роста",
        "sourceFileIds": ["source-notes"],
        "immutableElements": ["Фактические значения"],
        "locale": "ru-RU",
        "filename": filename,
        "acceptanceCriteria": [
            "Русский текст читается без подмены символов",
            "Фактологические слайды имеют источники",
            "Графики остаются редактируемыми",
        ],
    }


def _base_spec(filename: str = "presentation.pptx") -> dict:
    return {
        "job": _job(filename),
        "slides": [
            {
                "layout": "title",
                "title": "План роста на 2027 год",
                "subtitle": "Решение для совета директоров",
            },
            {
                "layout": "claim",
                "title": "Главный вывод",
                "claim": "Корпоративный сегмент обеспечит большую часть роста",
                "support": "План опирается на подтверждённую воронку и текущую экономику сделок.",
                "source": "Источник: финансовая модель, лист «План»",
            },
            {
                "layout": "chart",
                "title": "Выручка ускоряется после запуска корпоративного канала",
                "chart": {
                    "type": "column",
                    "categories": ["2025", "2026", "2027"],
                    "series": [{"name": "Выручка, млн ₽", "values": [82, 114, 161]}],
                },
                "takeaway": "+41% к 2026 году",
                "source": "Источник: финансовая модель, лист «План»",
            },
            {
                "layout": "summary",
                "title": "Что нужно утвердить сегодня",
                "bullets": ["Бюджет первой волны", "Контрольные точки", "Ежемесячный обзор воронки"],
            },
        ],
        "sources": [
            {"label": "Финансовая модель", "url": "source-notes — лист «План»"}
        ],
        "outputPdf": False,
    }


def _save(spec: dict, output: Path) -> Presentation:
    deck, _changes = BUILDER._build(spec, output)
    deck.save(output)
    return Presentation(output)


def _slide_xml(path: Path, number: int) -> bytes:
    with zipfile.ZipFile(path) as archive:
        return archive.read(f"ppt/slides/slide{number}.xml")


class PresentationBuilderTests(unittest.TestCase):
    def test_matrix_is_russian_first_and_covers_ten_distinct_scenarios(self):
        cases = json.loads((Path(__file__).with_name("cases.json")).read_text(encoding="utf-8"))

        self.assertEqual(len(cases), 10)
        self.assertEqual(len({case["id"] for case in cases}), 10)
        self.assertGreaterEqual(Counter(case["locale"] for case in cases)["ru-RU"], 7)
        self.assertTrue(all(case["prompt"] and case["requirements"] for case in cases))

    def test_default_deck_uses_cross_platform_cyrillic_font_and_no_ai_accent_stripes(self):
        with tempfile.TemporaryDirectory() as folder:
            output = Path(folder) / "presentation.pptx"
            deck = _save(_base_spec(), output)

        names = {shape.name for slide in deck.slides for shape in slide.shapes}
        fonts = {
            run.font.name
            for slide in deck.slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            for paragraph in shape.text_frame.paragraphs
            for run in paragraph.runs
            if run.text.strip()
        }
        self.assertNotIn("Title accent", names)
        self.assertNotIn("Claim accent", names)
        self.assertNotIn("Column accent", names)
        self.assertEqual(fonts, {"Arial"})

    def test_minimum_type_scale_is_enforced_for_default_layouts(self):
        with tempfile.TemporaryDirectory() as folder:
            output = Path(folder) / "presentation.pptx"
            deck = _save(_base_spec(), output)

        for slide in deck.slides:
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if not run.text.strip() or run.font.size is None:
                            continue
                        if shape.name in {"Source note", "Slide number", "Image caption"}:
                            continue
                        minimum = 35 if shape.name == "Slide title" else 16
                        self.assertGreaterEqual(run.font.size.pt, minimum, shape.name)

    def test_long_russian_slide_title_reserves_content_space(self):
        spec = _base_spec()
        spec["slides"] = [
            {
                "layout": "comparison",
                "title": "Перераспределение ответственности сокращает путь от подписания договора до первой измеримой ценности",
                "left": {"heading": "До изменения", "bullets": ["Один владелец"]},
                "right": {"heading": "После изменения", "bullets": ["Разделённая ответственность"]},
                "source": "Источник: анализ внедрений",
            }
        ]
        with tempfile.TemporaryDirectory() as folder:
            output = Path(folder) / "presentation.pptx"
            deck = _save(spec, output)
            title = next(shape for shape in deck.slides[0].shapes if shape.name == "Slide title")
            headings = [shape for shape in deck.slides[0].shapes if shape.name == "Column heading"]

        self.assertGreater(title.height, Inches(0.92))
        self.assertTrue(all(heading.top >= title.top + title.height for heading in headings))

    def test_metric_and_process_layouts_remain_editable_powerpoint_objects(self):
        spec = _base_spec()
        spec["slides"] = [
            {
                "layout": "title",
                "title": "Операционный план",
                "subtitle": "Ключевые показатели и этапы",
            },
            {
                "layout": "metrics",
                "title": "Экономика уже подтверждает масштабирование",
                "metrics": [
                    {"value": "41%", "label": "рост выручки"},
                    {"value": "4,2%", "label": "отток"},
                    {"value": "7 мес.", "label": "окупаемость"},
                ],
                "source": "Источник: управленческий отчёт",
            },
            {
                "layout": "process",
                "title": "Запуск проходит четыре контрольные точки",
                "steps": [
                    {"title": "Диагностика", "detail": "Неделя 1"},
                    {"title": "Пилот", "detail": "Недели 2–4"},
                    {"title": "Масштаб", "detail": "Месяцы 2–3"},
                    {"title": "Контроль", "detail": "Ежемесячно"},
                ],
                "source": "Источник: план запуска",
            },
        ]
        with tempfile.TemporaryDirectory() as folder:
            output = Path(folder) / "presentation.pptx"
            deck = _save(spec, output)

        names = {shape.name for slide in deck.slides for shape in slide.shapes}
        self.assertIn("Metric value", names)
        self.assertIn("Process step title", names)
        self.assertFalse(any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for slide in deck.slides for shape in slide.shapes))

    def test_native_chart_axis_ids_are_valid_unsigned_openxml_values(self):
        with tempfile.TemporaryDirectory() as folder:
            output = Path(folder) / "presentation.pptx"
            _save(_base_spec(), output)
            with zipfile.ZipFile(output) as archive:
                chart_xml = archive.read("ppt/charts/chart1.xml")

        self.assertNotRegex(chart_xml.decode("utf-8"), r'<c:(?:axId|crossAx) val="-')

    def test_decorative_shapes_do_not_inherit_default_theme_shadows(self):
        with tempfile.TemporaryDirectory() as folder:
            output = Path(folder) / "presentation.pptx"
            _save(_base_spec(), output)
            with zipfile.ZipFile(output) as archive:
                slide_xml = b"".join(
                    archive.read(name)
                    for name in archive.namelist()
                    if name.startswith("ppt/slides/slide") and name.endswith(".xml")
                )

        self.assertNotIn(b"<a:effectRef", slide_xml)

    def test_fact_slide_without_source_fails_structural_qa(self):
        spec = _base_spec()
        del spec["slides"][2]["source"]
        with tempfile.TemporaryDirectory() as folder:
            output = Path(folder) / "presentation.pptx"
            _save(spec, output)
            _checks, issues = BUILDER._check_structure(output, spec)

        self.assertTrue(any(issue["code"] == "missing-slide-source" and issue["severity"] == "critical" for issue in issues))

    def test_structural_qa_rejects_a_source_that_names_only_a_site(self):
        """A front-page link lets nobody check the fact it is attached to."""
        spec = _base_spec()
        spec["sources"] = [
            {"label": "CNBC — 27.07.2026", "url": "https://www.cnbc.com/2026/07/27/apple.html"},
            {"label": "Statista — 01.07.2026", "url": "https://www.statista.com/"},
        ]
        with tempfile.TemporaryDirectory() as folder:
            output = Path(folder) / "presentation.pptx"
            _save(spec, output)
            checks, issues = BUILDER._check_structure(output, spec)

        check = next(check for check in checks if check["name"] == "source-specificity")
        self.assertEqual(check["status"], "failed")
        self.assertIn("statista.com", check["message"])
        self.assertNotIn("cnbc.com", check["message"])
        self.assertTrue(
            any(
                issue["code"] == "unspecific-source" and issue["severity"] == "critical"
                for issue in issues
            )
        )

    def test_a_source_pointing_at_a_real_page_passes(self):
        spec = _base_spec()
        spec["sources"] = [
            {"label": "CNBC — 27.07.2026", "url": "https://www.cnbc.com/2026/07/27/apple.html"},
            # A query names a specific view of a site just as a path does; rejecting
            # it would refuse legitimate sources such as a sorted ranking table.
            {"label": "CompaniesMarketCap", "url": "https://companiesmarketcap.com/?sort=marketcap"},
            {"label": "Финансовая модель", "url": "model.xlsx — лист «План»"},
        ]
        with tempfile.TemporaryDirectory() as folder:
            output = Path(folder) / "presentation.pptx"
            _save(spec, output)
            checks, issues = BUILDER._check_structure(output, spec)

        check = next(check for check in checks if check["name"] == "source-specificity")
        self.assertEqual(check["status"], "passed")
        self.assertFalse(any(issue["code"] == "unspecific-source" for issue in issues))

    def test_sources_slide_shows_the_site_and_keeps_the_address_in_the_notes(self):
        """A raw article URL printed on a slide is unreadable and unclickable from
        a projector, so the slide carries the site and the notes carry the address."""
        url = "https://www.cnbc.com/2026/07/27/apple-most-valuable-company-nvidia.html"
        spec = _base_spec()
        spec["sources"] = [
            {"label": "CNBC — 27.07.2026", "url": url},
            {"label": "Финансовая модель", "url": "model.xlsx — лист «План»"},
        ]
        with tempfile.TemporaryDirectory() as folder:
            output = Path(folder) / "presentation.pptx"
            deck = _save(spec, output)

            sources_slide = deck.slides[-1]
            shown = " | ".join(
                shape.text_frame.text
                for shape in sources_slide.shapes
                if getattr(shape, "has_text_frame", False)
            )
            self.assertIn("cnbc.com", shown)
            self.assertNotIn(url, shown)
            self.assertIn("model.xlsx — лист «План»", shown)

            notes = sources_slide.notes_slide.notes_text_frame.text
            self.assertIn(url, notes)
            self.assertNotIn("model.xlsx", notes)

    def test_template_sources_slide_also_keeps_the_address_in_the_notes(self):
        """Template decks render through their own path; the first fix reached only
        the builder-native one and left raw URLs on inherited layouts."""
        url = "https://www.cnbc.com/2026/07/27/apple-most-valuable-company-nvidia.html"
        with tempfile.TemporaryDirectory() as folder:
            root = Path(folder)
            template_path = root / "template.pptx"
            template = Presentation()
            layout = next(
                item for item in template.slide_layouts if "title and content" in item.name.lower()
            )
            template.save(template_path)
            entries = [{"label": "CNBC — 27.07.2026", "url": url}]
            spec = {
                "job": _job("output.pptx"),
                "templatePath": str(template_path),
                "slides": [
                    {
                        "layout": "sources",
                        "templateLayout": layout.name,
                        "title": "Источники",
                        "entries": entries,
                    }
                ],
                "sources": entries,
            }
            output = root / "output.pptx"
            deck, _changes = BUILDER._build(spec, output)
            deck.save(str(output))

            slide = Presentation(str(output)).slides[-1]
            shown = " | ".join(
                shape.text_frame.text
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
            )
            self.assertIn("cnbc.com", shown)
            self.assertNotIn(url, shown)
            self.assertIn(url, slide.notes_slide.notes_text_frame.text)

    def test_slide_title_shrinks_instead_of_growing_into_the_content(self):
        """Content is placed below the height reserved for the title, and that
        height comes from a character count that cannot be right for both narrow
        and wide letters. The renderer must shrink the title to the box it was
        promised rather than wrap into the table underneath."""
        from pptx.enum.text import MSO_AUTO_SIZE

        spec = _base_spec()
        spec["slides"] = [
            {
                "layout": "bullets",
                "title": "Знакомьтесь, виновники инфляции ваших нервов",
                "bullets": ["Первое", "Второе"],
            }
        ]
        with tempfile.TemporaryDirectory() as folder:
            output = Path(folder) / "presentation.pptx"
            deck = _save(spec, output)

        title = next(
            shape
            for shape in deck.slides[0].shapes
            if getattr(shape, "has_text_frame", False)
            and shape.text_frame.text.startswith("Знакомьтесь")
        )
        self.assertEqual(title.text_frame.auto_size, MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE)

    def test_sources_slide_fits_a_real_label_without_overflowing(self):
        """The first live deck failed QA here: a 42-character label in a 4-inch
        column wrapped onto the divider, and the model burned its whole turn
        trying to satisfy a check the layout made unsatisfiable."""
        spec = _base_spec()
        spec["sources"] = [
            {
                "label": "Capital.com — ИИ-компании по капитализации",
                "url": "https://capital.com/en-int/markets/shares/largest-ai-companies",
            },
            {
                "label": "AlphaSense — крупнейшие компании по капитализации",
                "url": "https://www.alpha-sense.com/largest-companies-by-market-cap/",
            },
        ]
        with tempfile.TemporaryDirectory() as folder:
            output = Path(folder) / "presentation.pptx"
            _save(spec, output)
            _checks, issues = BUILDER._check_structure(output, spec)

        overflow = [
            issue
            for issue in issues
            if issue["code"] == "text-overflow-risk" and "Source" in str(issue.get("message", ""))
        ]
        self.assertEqual(overflow, [])

    def test_sources_rows_keep_the_divider_below_the_label(self):
        spec = _base_spec()
        spec["sources"] = [
            {"label": "Довольно длинная подпись источника на две строки", "url": "https://example.com/a"},
            {"label": "Вторая подпись", "url": "https://example.com/b"},
        ]
        with tempfile.TemporaryDirectory() as folder:
            output = Path(folder) / "presentation.pptx"
            deck = _save(spec, output)

        slide = deck.slides[-1]
        label = next(
            shape for shape in slide.shapes if shape.name == "Source label"
        )
        divider = next(shape for shape in slide.shapes if shape.name == "Source divider")
        self.assertGreaterEqual(
            divider.top,
            label.top + label.height,
            "the divider is drawn through the label it should sit under",
        )

    def test_a_file_source_still_fits_its_own_location(self):
        """A file location ("management-report.xlsx — июль 2026") is a sentence,
        not a domain. Taller rows must not come at the cost of its width."""
        spec = _base_spec()
        spec["sources"] = [
            {"label": "Управленческий отчёт", "url": "management-report.xlsx — июль 2026"},
            {"label": "CRM-выгрузка", "url": "crm-export.csv — срез от 31.07.2026"},
        ]
        with tempfile.TemporaryDirectory() as folder:
            output = Path(folder) / "presentation.pptx"
            deck = _save(spec, output)
            _checks, issues = BUILDER._check_structure(output, spec)

        slide = deck.slides[-1]
        location = next(shape for shape in slide.shapes if shape.name == "Source location")
        divider = next(shape for shape in slide.shapes if shape.name == "Source divider")
        self.assertGreaterEqual(divider.top, location.top + location.height)
        self.assertEqual(
            [issue for issue in issues if issue["code"] == "text-overflow-risk"], []
        )

    def test_sources_rows_stay_tall_enough_at_the_eight_entry_limit(self):
        """Rows share the free height, but never shrink below what a two-line
        label needs — at the maximum of eight entries the share alone would."""
        spec = _base_spec()
        spec["sources"] = [
            {
                "label": f"Источник {index} — довольно длинная подпись на две строки точно",
                "url": f"https://example.com/page-{index}",
            }
            for index in range(1, 9)
        ]
        with tempfile.TemporaryDirectory() as folder:
            output = Path(folder) / "presentation.pptx"
            deck = _save(spec, output)
            _checks, issues = BUILDER._check_structure(output, spec)

        slide = deck.slides[-1]
        labels = [shape for shape in slide.shapes if shape.name == "Source label"]
        dividers = [shape for shape in slide.shapes if shape.name == "Source divider"]
        self.assertEqual(len(labels), 8)
        for label, divider in zip(labels, dividers):
            self.assertGreaterEqual(divider.top, label.top + label.height)
        self.assertEqual(
            [issue for issue in issues if issue["code"] == "text-overflow-risk"], []
        )
        self.assertLess(
            (labels[-1].top + labels[-1].height) / 914400,
            7.03,
            "the last label runs into the source note at the foot of the slide",
        )

    def test_claim_marker_shares_the_centre_line_of_the_claim_it_marks(self):
        """The claim is centred in its box; a marker pinned near the top of that box
        drifts away from the words — by an inch on a two-line claim."""
        spec = _base_spec()
        spec["slides"] = [
            {
                "layout": "claim",
                "title": "Ключевое ограничение находится не в продажах",
                "claim": "Рост зависит от первых четырёх недель работы клиента",
                "support": "Поэтому ресурсы нужно перенести в раннее внедрение.",
                "source": "Источник: управленческий отчёт",
            }
        ]
        with tempfile.TemporaryDirectory() as folder:
            output = Path(folder) / "presentation.pptx"
            deck = _save(spec, output)

        slide = deck.slides[0]
        # The slide's background is an auto shape too; the marker is the small one.
        marker = next(
            shape
            for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.width == Inches(0.2)
        )
        claim = next(
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and "Рост зависит" in shape.text_frame.text
        )
        marker_centre = marker.top + marker.height / 2
        claim_centre = claim.top + claim.height / 2
        self.assertLess(
            abs(marker_centre - claim_centre),
            Inches(0.02),
            "the marker must sit on the centre line the claim text is anchored to",
        )

    def test_a_title_keeps_a_readable_break_above_its_content(self):
        """Below half the title's own line height the heading reads as part of the
        text under it; the builder had 0.42 of it and a three-line title sat on its
        own bullets."""
        gap = BUILDER.TITLE_GAP / BUILDER.TITLE_LINE_HEIGHT
        self.assertGreaterEqual(gap, 0.5)
        self.assertLessEqual(gap, 0.75)

        content_top = {}
        for title in ("Короткий заголовок", "Перераспределение ответственности сокращает путь от подписания договора до первой ценности"):
            spec = _base_spec()
            spec["slides"] = [{"layout": "bullets", "title": title, "bullets": ["Пункт"]}]
            with tempfile.TemporaryDirectory() as folder:
                output = Path(folder) / "presentation.pptx"
                deck = _save(spec, output)
            slide = deck.slides[0]
            heading = next(
                shape
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and title in shape.text_frame.text
            )
            body = next(
                shape
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and "Пункт" in shape.text_frame.text
            )
            content_top[title] = body.top - (heading.top + heading.height)

        for title, measured in content_top.items():
            self.assertGreaterEqual(
                measured,
                Inches(BUILDER.TITLE_LINE_HEIGHT * 0.5),
                f"content sits too close under the title: {title[:40]}",
            )

    def test_structural_qa_rejects_text_exceeding_shape_capacity(self):
        spec = _base_spec()
        spec["slides"] = [
            {
                "layout": "bullets",
                "title": "Слишком плотный слайд",
                "bullets": ["Очень длинный тезис " * 35],
            }
        ]
        with tempfile.TemporaryDirectory() as folder:
            output = Path(folder) / "presentation.pptx"
            _save(spec, output)
            _checks, issues = BUILDER._check_structure(output, spec)

        overflow = next(
            issue for issue in issues if issue["code"] == "text-overflow-risk"
        )
        self.assertEqual(overflow["severity"], "critical")
        self.assertEqual(overflow["target"], "Slide 1")
        # Naming the shape alone gave two identical messages for two long labels
        # on one slide; the reader could not tell which text to shorten.
        self.assertIn("Очень длинный тезис", overflow["message"])

    def test_targeted_edit_preserves_unaffected_slide_xml(self):
        with tempfile.TemporaryDirectory() as folder:
            root = Path(folder)
            source = root / "source.pptx"
            _save(_base_spec(source.name), source)
            source_bytes = source.read_bytes()
            source_properties = Presentation(source).core_properties
            source_title = source_properties.title
            source_subject = source_properties.subject
            source_comments = source_properties.comments
            before = {number: _slide_xml(source, number) for number in (1, 2, 4)}
            output = root / "updated.pptx"
            revision_job = _job(output.name)
            revision_job["goal"] = "A different request description must not rewrite deck metadata"
            revision_job["audience"] = "A different audience must not rewrite deck metadata"
            spec = {
                "job": revision_job,
                "inputPath": str(source),
                "edits": [
                    {
                        "slide": 3,
                        "replacements": {"+41% к 2026 году": "+44% к 2026 году"},
                        "summary": "Уточнён прогноз роста",
                    }
                ],
            }
            _save(spec, output)

            self.assertEqual(source.read_bytes(), source_bytes)
            for number, xml in before.items():
                self.assertEqual(_slide_xml(output, number), xml, f"slide {number}")
            self.assertEqual(
                BUILDER._changed_package_parts(source, output),
                {"ppt/slides/slide3.xml"},
            )
            output_properties = Presentation(output).core_properties
            self.assertEqual(output_properties.title, source_title)
            self.assertEqual(output_properties.subject, source_subject)
            self.assertEqual(output_properties.comments, source_comments)

    def test_targeted_edit_qa_rejects_an_unrequested_package_change(self):
        with tempfile.TemporaryDirectory() as folder:
            root = Path(folder)
            source = root / "source.pptx"
            _save(_base_spec(source.name), source)
            output = root / "updated.pptx"
            spec = {
                "job": _job(output.name),
                "inputPath": str(source),
                "edits": [
                    {
                        "slide": 3,
                        "replacements": {"+41% к 2026 году": "+44% к 2026 году"},
                    }
                ],
            }
            _save(spec, output)
            damaged = Presentation(output)
            title = next(
                shape
                for shape in damaged.slides[1].shapes
                if getattr(shape, "has_text_frame", False) and shape.text.strip()
            )
            title.text_frame.paragraphs[0].runs[0].text = "Незапрошенное изменение"
            damaged.save(output)

            checks, issues = BUILDER._check_structure(output, spec)

        scope_check = next(check for check in checks if check["name"] == "targeted-edit-scope")
        self.assertEqual(scope_check["status"], "failed")
        self.assertTrue(
            any(
                issue["code"] == "edit-scope"
                and issue["severity"] == "critical"
                and "outside" in issue["message"]
                for issue in issues
            )
        )

    def test_targeted_edit_qa_fails_closed_when_package_comparison_errors(self):
        with tempfile.TemporaryDirectory() as folder:
            root = Path(folder)
            source = root / "source.pptx"
            _save(_base_spec(source.name), source)
            output = root / "updated.pptx"
            spec = {
                "job": _job(output.name),
                "inputPath": str(source),
                "edits": [
                    {
                        "slide": 3,
                        "replacements": {"+41% к 2026 году": "+44% к 2026 году"},
                    }
                ],
            }
            _save(spec, output)

            with mock.patch.object(
                BUILDER,
                "_changed_package_parts",
                side_effect=zipfile.BadZipFile("invalid package"),
            ):
                checks, issues = BUILDER._check_structure(output, spec)

        scope_check = next(check for check in checks if check["name"] == "targeted-edit-scope")
        self.assertEqual(scope_check["status"], "failed")
        self.assertIn("BadZipFile", scope_check["message"])
        self.assertTrue(
            any(
                issue["code"] == "edit-scope" and issue["severity"] == "critical"
                for issue in issues
            )
        )

    def test_targeted_edit_qa_refuses_to_vouch_for_a_duplicate_part_package(self):
        with tempfile.TemporaryDirectory() as folder:
            root = Path(folder)
            source = root / "source.pptx"
            _save(_base_spec(source.name), source)
            duplicated = root / "duplicated.pptx"
            with zipfile.ZipFile(source) as original, zipfile.ZipFile(duplicated, "w") as clone:
                for info in original.infolist():
                    clone.writestr(info, original.read(info.filename))
                clone.writestr("ppt/slides/slide3.xml", original.read("ppt/slides/slide3.xml"))
            output = root / "updated.pptx"
            spec = {
                "job": _job(output.name),
                "inputPath": str(duplicated),
                "edits": [
                    {
                        "slide": 3,
                        "replacements": {"+41% к 2026 году": "+44% к 2026 году"},
                    }
                ],
            }
            _save(spec, output)

            checks, issues = BUILDER._check_structure(output, spec)

        scope_check = next(check for check in checks if check["name"] == "targeted-edit-scope")
        self.assertEqual(scope_check["status"], "failed")
        self.assertIn("ValueError", scope_check["message"])
        self.assertTrue(
            any(
                issue["code"] == "edit-scope" and issue["severity"] == "critical"
                for issue in issues
            )
        )

    def test_skill_runtime_contract_keeps_scratch_and_final_files_separate(self):
        instructions = (ROOT / "skill/pptx/SKILL.md").read_text(encoding="utf-8")

        self.assertIn("Treat the builder as an executable", instructions)
        self.assertIn("Never reconstruct an attached presentation", instructions)
        self.assertIn("same-call scratch", instructions)
        self.assertIn("Final user files must be direct children of `/mnt/data`", instructions)
        self.assertIn("Do not deliver or mention the JSON spec", instructions)

    def test_template_mode_rejects_layout_without_inherited_content_slots(self):
        with tempfile.TemporaryDirectory() as folder:
            root = Path(folder)
            template_path = root / "template.pptx"
            template = Presentation()
            blank_layout = next(
                layout
                for layout in template.slide_layouts
                if layout.name.lower() == "blank"
            )
            template.save(template_path)
            spec = {
                "job": _job("output.pptx"),
                "templatePath": str(template_path),
                "slides": [
                    {
                        "layout": "claim",
                        "templateLayout": blank_layout.name,
                        "title": "Нельзя рисовать поверх пустого шаблона",
                        "claim": "Нужны унаследованные слоты",
                    }
                ],
            }

            with self.assertRaisesRegex(ValueError, "inherited content placeholder"):
                BUILDER._build(spec, root / "output.pptx")

    def test_template_mode_enables_native_text_fitting_for_inherited_placeholders(self):
        with tempfile.TemporaryDirectory() as folder:
            root = Path(folder)
            template_path = root / "template.pptx"
            template = Presentation()
            title_and_content = template.slide_layouts[1]
            template.save(template_path)
            output = root / "output.pptx"
            spec = {
                "job": _job(output.name),
                "templatePath": str(template_path),
                "slides": [
                    {
                        "layout": "bullets",
                        "templateLayout": title_and_content.name,
                        "title": "Длинный русский заголовок должен оставаться внутри унаследованного поля",
                        "bullets": [
                            "Первый содержательный тезис для проверки заполнения",
                            "Второй содержательный тезис для проверки заполнения",
                        ],
                    },
                    {
                        "layout": "bullets",
                        "templateLayout": title_and_content.name,
                        "title": "Короткий заголовок",
                        "bullets": ["Наследует размер шрифта шаблона"],
                    },
                ],
            }

            deck = _save(spec, output)
            slide = deck.slides[0]
            title = slide.shapes.title
            bodies = BUILDER._template_body_placeholders(slide)

            self.assertEqual(title.text_frame.auto_size, MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE)
            self.assertEqual(bodies[0].text_frame.auto_size, MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE)
            self.assertEqual(title.text_frame.vertical_anchor, MSO_ANCHOR.TOP)
            title_sizes = [
                run.font.size.pt
                for paragraph in title.text_frame.paragraphs
                for run in paragraph.runs
                if run.font.size is not None
            ]
            self.assertTrue(title_sizes)
            self.assertLessEqual(max(title_sizes), 30)
            short_title = deck.slides[1].shapes.title
            self.assertTrue(
                all(
                    run.font.size is None
                    for paragraph in short_title.text_frame.paragraphs
                    for run in paragraph.runs
                )
            )


if __name__ == "__main__":
    unittest.main()
